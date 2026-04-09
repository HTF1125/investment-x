from __future__ import annotations

import json
import os
import textwrap
from io import BytesIO
from typing import Optional
from datetime import datetime

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import yfinance as yf
from fastapi import APIRouter, Depends, HTTPException, Query, Body

from ix.api.dependencies import get_current_user
from fastapi.responses import StreamingResponse
from plotly.subplots import make_subplots
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib import utils
from pptx import Presentation
from pptx.util import Inches, Pt

from ix.common import get_logger
from ix.core.technical.ohlcv_indicators import (
    _normalize_yf,
    _compute_rsi,
    _compute_squeeze_momentum,
    _compute_supertrend,
    _compute_moving_averages,
    _compute_bollinger_bands,
    _compute_vwap,
    _compute_stochastic,
    _compute_atr,
    _fmt_price,
    _find_support_resistance,
    _find_swing_points,
    _fit_trendline,
)
from ix.core.technical.chart_builder import (
    _add_fib_zone,
    _build_figure,
    _render_chart_to_image,
    _clean_markdown,
)

router = APIRouter()
logger = get_logger(__name__)



@router.get("/technical/elliott")
def technical_elliott(
    ticker: str = Query("SPY"),
    period: str = Query("2y"),
    interval: str = Query("1d"),
    start: Optional[str] = Query(None),
    end: Optional[str] = Query(None),
    setup_from: int = Query(9, ge=1, le=9),
    countdown_from: int = Query(13, ge=1, le=13),
    label_cooldown: int = Query(0, ge=0, le=20),
    show_macd: bool = Query(True),
    show_rsi: bool = Query(True),
    show_sqz: bool = Query(False),
    sqz_bb_len: int = Query(20, ge=5, le=50),
    sqz_bb_mult: float = Query(2.0, ge=0.5, le=5.0),
    sqz_kc_len: int = Query(20, ge=5, le=50),
    sqz_kc_mult: float = Query(1.5, ge=0.5, le=5.0),
    _user=Depends(get_current_user),
):
    try:
        tk = ticker.strip().upper()
        raw = yf.download(tk, period=period, interval=interval, auto_adjust=False, progress=False)
        if raw is None or raw.empty:
            raise HTTPException(status_code=404, detail=f"No data for ticker '{ticker}'.")
        df = _normalize_yf(raw, tk)
        if df.empty:
            raise HTTPException(status_code=404, detail=f"No OHLC data for ticker '{ticker}'.")
        vis_start = pd.to_datetime(start) if start else None
        vis_end = pd.to_datetime(end) if end else None
        fig = _build_figure(
            df,
            tk,
            setup_from,
            countdown_from,
            label_cooldown,
            visible_start=vis_start,
            visible_end=vis_end,
            show_macd=show_macd,
            show_rsi=show_rsi,
            show_sqz=show_sqz,
            sqz_bb_len=sqz_bb_len,
            sqz_bb_mult=sqz_bb_mult,
            sqz_kc_len=sqz_kc_len,
            sqz_kc_mult=sqz_kc_mult,
        )
        return json.loads(pio.to_json(fig, engine="json"))
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail="Failed to build technical chart")


@router.get("/technical/overlays")
def technical_overlays(
    ticker: str = Query("SPY"),
    interval: str = Query("1d"),
    # Squeeze Momentum params
    sqz: bool = Query(False),
    sqz_bb_len: int = Query(20, ge=5, le=50),
    sqz_bb_mult: float = Query(2.0, ge=0.5, le=5.0),
    sqz_kc_len: int = Query(20, ge=5, le=50),
    sqz_kc_mult: float = Query(1.5, ge=0.5, le=5.0),
    # Supertrend params
    st: bool = Query(False),
    st_period: int = Query(10, ge=2, le=50),
    st_mult: float = Query(3.0, ge=0.5, le=10.0),
    # Moving averages — comma-separated list of "TYPE:period:color"
    # e.g. "SMA:20:#f59e0b,EMA:50:#38bdf8"
    mas: Optional[str] = Query(None),
    # Bollinger Bands
    bb: bool = Query(False),
    bb_len: int = Query(20, ge=5, le=100),
    bb_mult: float = Query(2.0, ge=0.5, le=5.0),
    # VWAP
    vwap: bool = Query(False),
    # Stochastic Oscillator
    stoch: bool = Query(False),
    stoch_k: int = Query(14, ge=2, le=50),
    stoch_d: int = Query(3, ge=1, le=20),
    stoch_smooth: int = Query(3, ge=1, le=10),
    # ATR
    atr: bool = Query(False),
    atr_period: int = Query(14, ge=2, le=50),
    _user=Depends(get_current_user),
):
    """Return overlay data (Squeeze Momentum, Supertrend, extra MAs) as raw series."""
    try:
        tk = ticker.strip().upper()
        period_map = {"1d": "2y", "1wk": "5y", "1mo": "10y"}
        yf_period = period_map.get(interval, "2y")
        raw = yf.download(tk, period=yf_period, interval=interval, auto_adjust=False, progress=False)
        if raw is None or raw.empty:
            raise HTTPException(status_code=404, detail=f"No data for ticker '{tk}'.")
        df = _normalize_yf(raw, tk)
        if df.empty:
            raise HTTPException(status_code=404, detail=f"No OHLC data for ticker '{tk}'.")

        dates = [d.isoformat() if hasattr(d, "isoformat") else str(d) for d in df.index]
        result: dict = {"dates": dates}

        if sqz:
            sqz_df = _compute_squeeze_momentum(
                df,
                bb_length=sqz_bb_len,
                bb_mult=sqz_bb_mult,
                kc_length=sqz_kc_len,
                kc_mult=sqz_kc_mult,
            )
            result["squeeze"] = {
                "val": [None if np.isnan(v) else round(float(v), 6) for v in sqz_df["val"]],
                "bar_color": sqz_df["bar_color"].tolist(),
                "sqz_on": sqz_df["sqz_on"].tolist(),
                "sqz_off": sqz_df["sqz_off"].tolist(),
                "no_sqz": sqz_df["no_sqz"].tolist(),
                "sqz_dot_color": sqz_df["sqz_dot_color"].tolist(),
            }

        if st:
            st_df = _compute_supertrend(df, atr_period=st_period, multiplier=st_mult)
            result["supertrend"] = {
                "trend": st_df["trend"].tolist(),
                "up": [None if np.isnan(v) else round(float(v), 4) for v in st_df["up"]],
                "dn": [None if np.isnan(v) else round(float(v), 4) for v in st_df["dn"]],
                "buy": st_df["buy"].tolist(),
                "sell": st_df["sell"].tolist(),
            }

        if mas:
            ma_configs = []
            for part in mas.split(","):
                parts = part.strip().split(":")
                if len(parts) >= 2:
                    ma_configs.append({
                        "type": parts[0].upper(),
                        "period": int(parts[1]),
                        "color": parts[2] if len(parts) > 2 else "#94a3b8",
                    })
            if ma_configs:
                result["moving_averages"] = _compute_moving_averages(df, ma_configs)

        if bb:
            result["bollinger"] = _compute_bollinger_bands(df, length=bb_len, mult=bb_mult)

        if vwap:
            result["vwap"] = _compute_vwap(df)

        if stoch:
            result["stochastic"] = _compute_stochastic(
                df, k_period=stoch_k, d_period=stoch_d, smooth_k=stoch_smooth
            )

        if atr:
            result["atr"] = _compute_atr(df, period=atr_period)

        return result
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail="Failed to compute overlays")


@router.post("/technical/export")
def export_technical_report(
    ticker: str = Query(...),
    format: str = Query("pdf"), # "pdf" or "pptx"
    summary: str = Body(..., embed=True),
    # Chart params to reproduce the image
    interval: str = Query("1d"),
    setup_from: int = Query(9),
    countdown_from: int = Query(13),
    label_cooldown: int = Query(0),
    show_macd: bool = Query(True),
    show_rsi: bool = Query(True),
    _user=Depends(get_current_user),
):
    try:
        tk = ticker.strip().upper()
        # 1. Generate the chart image
        raw = yf.download(tk, period="2y", interval=interval, auto_adjust=False, progress=False)
        if raw is None or raw.empty:
            raise HTTPException(status_code=404, detail=f"No data for ticker '{ticker}'.")
        df = _normalize_yf(raw, tk)
        
        fig = _build_figure(
            df, tk, setup_from, countdown_from, label_cooldown,
            show_macd=show_macd, show_rsi=show_rsi
        )
        # Apply export-friendly theme
        fig.update_layout(
            template="plotly_white",
            paper_bgcolor="white",
            plot_bgcolor="#fcfcfc",
            font=dict(color="black", size=14),
            margin=dict(l=80, r=40, t=100, b=80),
        )
        chart_img = _render_chart_to_image(fig)

        filename = f"InvestmentX_{tk}_Analysis_{datetime.now().strftime('%Y%m%d')}.{format}"

        if format.lower() == "pptx":
            prs = Presentation()
            
            # Slide 1: Title
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"Technical Analysis: {tk}"
            subtitle.text = f"Generated by Investment-X Engine\nDate: {datetime.now().strftime('%Y-%m-%d')}"

            # Slide 2: Chart
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(chart_img, Inches(0.5), Inches(1), width=Inches(9))
            
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = f"{tk} Price Action & Indicators ({interval})"
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(24)

            # Slide 3+: Intelligence Report (Split if too long)
            content_lines = _clean_markdown(summary).split("\n")
            chunks = [content_lines[i:i + 15] for i in range(0, len(content_lines), 15)]
            
            for i, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Intelligence Report {' (cont.)' if i > 0 else ''}"
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                tf.text = "\n".join(chunk)
                for p in tf.paragraphs:
                    p.font.size = Pt(14)

            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )

        else: # Default to PDF
            buffer = BytesIO()
            c = canvas.Canvas(buffer, pagesize=letter)
            width, height = letter

            # Header
            c.setFont("Helvetica-Bold", 24)
            c.drawString(50, height - 60, f"Technical Analysis: {tk}")
            c.setFont("Helvetica", 12)
            c.setStrokeColorRGB(0.2, 0.5, 0.8)
            c.line(50, height - 75, width - 50, height - 75)
            
            c.drawString(50, height - 95, f"Ticker: {tk} | Interval: {interval} | Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

            # Chart Image (Centered)
            img_reader = utils.ImageReader(chart_img)
            c.drawImage(img_reader, 50, height - 420, width=width-100, height=300, preserveAspectRatio=True)

            # Report Text
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 450, "Intelligence Report")
            c.line(50, height - 455, 200, height - 455)

            c.setFont("Helvetica", 10)
            text_object = c.beginText(50, height - 480)
            text_object.setLeading(14)
            
            wrapped_text = ""
            for line in _clean_markdown(summary).split("\n"):
                wrapped_text += "\n".join(textwrap.wrap(line, width=95)) + "\n"
            
            for line in wrapped_text.split("\n"):
                if text_object.getY() < 50:
                    c.drawText(text_object)
                    c.showPage()
                    text_object = c.beginText(50, height - 50)
                    text_object.setFont("Helvetica", 10)
                    text_object.setLeading(14)
                text_object.textLine(line)
            
            c.drawText(text_object)
            c.showPage()
            c.save()
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )

    except Exception as e:
        logger.exception(
            "Technical report export failed for ticker=%s format=%s interval=%s: %s",
            ticker,
            format,
            interval,
            e,
        )
        raise HTTPException(status_code=500, detail="Export failed")


# ── Research Chart Endpoint ───────────────────────────────────────────────


@router.get("/technical/research-chart")
def get_research_chart(
    ticker: str = Query("SPY"),
    interval: str = Query("1d"),
    period: str = Query("1y"),
    _user=Depends(get_current_user),
):
    """Research-style chart: MACD + Stoch on top, candlestick + MAs in center, volume at bottom.

    Always fetches max available history so indicators (especially 200MA)
    are calculated on the full dataset.  The ``period`` parameter only
    controls the **initial visible x-axis range** — the user can pan left
    to see older data.
    """
    try:
        tk = ticker.strip().upper()

        # Fetch extra history so the longest MA (200) is warmed up
        # before the visible window, and users can pan left a bit.
        # We fetch 3× the requested period to give buffer, capped at 10y.
        period_years = {"1y": 1, "3y": 3, "5y": 5, "10y": 10}
        view_years = period_years.get(period, 1)
        fetch_years = min(view_years * 3, 10)
        fetch_period = f"{fetch_years}y"

        df = yf.download(tk, period=fetch_period, interval=interval, progress=False)
        if df.empty:
            raise HTTPException(status_code=404, detail=f"No data for {tk}")

        if isinstance(df.columns, pd.MultiIndex):
            df.columns = df.columns.get_level_values(0)
        df = df.reset_index()
        date_col = "Date" if "Date" in df.columns else "Datetime"
        dates = df[date_col].tolist()
        close = df["Close"].astype(float)
        high = df["High"].astype(float)
        low = df["Low"].astype(float)
        opn = df["Open"].astype(float)

        # Initial visible x-range = the requested period
        last_date = pd.Timestamp(dates[-1])
        view_start = last_date - pd.DateOffset(years=view_years)
        if view_start < pd.Timestamp(dates[0]):
            view_start = pd.Timestamp(dates[0])
        x_range_initial = [str(view_start.date()), str(last_date.date())]

        # ── Indicators ──────────────────────────────────────────────────

        ma_config = [(20, "#3b82f6"), (60, "#f59e0b"), (120, "#22c55e"), (200, "#a855f7")]
        ma_series = {}
        for p, _ in ma_config:
            if len(df) >= p:
                ma_series[p] = close.rolling(window=p).mean()

        ema12 = close.ewm(span=12, adjust=False).mean()
        ema26 = close.ewm(span=26, adjust=False).mean()
        macd_line = ema12 - ema26
        macd_signal = macd_line.ewm(span=9, adjust=False).mean()
        macd_hist = macd_line - macd_signal

        stoch_n = 14
        lowest_low = low.rolling(window=stoch_n).min()
        highest_high = high.rolling(window=stoch_n).max()
        fast_k = ((close - lowest_low) / (highest_high - lowest_low)) * 100
        slow_k = fast_k.rolling(window=3).mean()
        slow_d = slow_k.rolling(window=3).mean()

        # ── Layout: Row1=MACD, Row2=Stoch, Row3=Candlestick, Row4=Volume ─

        fig = make_subplots(
            rows=4, cols=1, shared_xaxes=True,
            row_heights=[0.13, 0.13, 0.62, 0.12],
            vertical_spacing=0.008,
        )

        # ── Row 1: MACD ─────────────────────────────────────────────────

        # Histogram with momentum-aware gradient colors
        hist_vals = macd_hist.tolist()
        hist_colors = []
        for i, v in enumerate(hist_vals):
            prev = hist_vals[i - 1] if i > 0 else 0
            if v >= 0:
                hist_colors.append("#34d399" if v >= prev else "#6ee7b7")  # bright/dim green
            else:
                hist_colors.append("#f87171" if v <= prev else "#fca5a5")  # bright/dim red

        fig.add_trace(go.Bar(
            x=dates, y=macd_hist, marker_color=hist_colors,
            marker_line_width=0, name="Hist", showlegend=False,
            hovertemplate="Hist: %{y:.4f}<extra></extra>",
        ), row=1, col=1)
        fig.add_trace(go.Scatter(
            x=dates, y=macd_line, mode="lines",
            line=dict(color="#38bdf8", width=1.3),
            name="MACD", showlegend=False,
            hovertemplate="MACD: %{y:.4f}<extra></extra>",
        ), row=1, col=1)
        fig.add_trace(go.Scatter(
            x=dates, y=macd_signal, mode="lines",
            line=dict(color="#fb923c", width=1, dash="dot"),
            name="Signal", showlegend=False,
            hovertemplate="Signal: %{y:.4f}<extra></extra>",
        ), row=1, col=1)
        fig.add_hline(y=0, line_color="rgba(148,163,184,0.15)", line_width=0.5, row=1, col=1)

        # ── Row 2: Stochastic ───────────────────────────────────────────

        # Fill between K and D for visual
        fig.add_trace(go.Scatter(
            x=dates, y=slow_k, mode="lines",
            line=dict(color="#3b82f6", width=1.3),
            name="Slow %K", showlegend=False,
            hovertemplate="Slow %%K: %{y:.1f}<extra></extra>",
        ), row=2, col=1)
        fig.add_trace(go.Scatter(
            x=dates, y=slow_d, mode="lines",
            line=dict(color="#f43f5e", width=1, dash="dot"),
            fill="tonexty", fillcolor="rgba(59,130,246,0.05)",
            name="Slow %D", showlegend=False,
            hovertemplate="Slow %%D: %{y:.1f}<extra></extra>",
        ), row=2, col=1)
        for ref_y in [20, 80]:
            fig.add_hline(
                y=ref_y, line_dash="dot",
                line_color="rgba(148,163,184,0.2)", line_width=0.5,
                row=2, col=1,
            )
        fig.update_yaxes(range=[0, 100], row=2, col=1)

        # ── Row 3: Candlestick + MAs ────────────────────────────────────

        fig.add_trace(go.Candlestick(
            x=dates, open=opn, high=high, low=low, close=close,
            increasing_line_color="#26a69a", increasing_fillcolor="#26a69a",
            decreasing_line_color="#ef5350", decreasing_fillcolor="#ef5350",
            name="OHLC", showlegend=False,
        ), row=3, col=1)

        annotations = []
        for ma_period, color in ma_config:
            if ma_period not in ma_series:
                continue
            mv = ma_series[ma_period]
            fig.add_trace(go.Scatter(
                x=dates, y=mv, mode="lines",
                line=dict(color=color, width=1.2),
                name=f"{ma_period}MA", showlegend=False,
                hovertemplate=f"{ma_period}MA: %{{y:.2f}}<extra></extra>",
            ), row=3, col=1)
            last_val = mv.dropna().iloc[-1] if not mv.dropna().empty else None
            if last_val is not None:
                annotations.append(dict(
                    x=1.0, y=float(last_val), xref="paper", yref="y3",
                    text=f"  {ma_period} {_fmt_price(float(last_val))}",
                    showarrow=False, xanchor="left",
                    font=dict(color=color, size=9, family="'Inter', sans-serif"),
                    bgcolor="rgba(0,0,0,0)",
                ))

        # Trendlines (on row 3 candlestick)
        swing_window = max(10, len(df) // 25)
        swing_hi_idx, swing_lo_idx = [], []
        start_i = max(swing_window, len(df) // 3)
        for i in range(start_i, len(df) - swing_window):
            seg_h = high.iloc[i - swing_window : i + swing_window + 1]
            if high.iloc[i] == seg_h.max():
                swing_hi_idx.append(i)
            seg_l = low.iloc[i - swing_window : i + swing_window + 1]
            if low.iloc[i] == seg_l.min():
                swing_lo_idx.append(i)

        for is_high, idxs, vals in [(True, swing_hi_idx, high), (False, swing_lo_idx, low)]:
            tl = _fit_trendline(dates, vals, idxs, is_high)
            if tl:
                fig.add_shape(
                    type="line",
                    x0=tl["x0"], y0=tl["y0"], x1=tl["x1"], y1=tl["y1"],
                    line=dict(color=tl["color"], width=1.5, dash="dash"),
                    row=3, col=1,
                )

        # Support/Resistance
        key_levels = _find_support_resistance(high, low, window=20, n_levels=5)
        for level in key_levels:
            fig.add_hline(
                y=level, line_dash="dot", line_color="rgba(148,163,184,0.25)",
                line_width=0.8, row=3, col=1,
            )
            annotations.append(dict(
                x=0.0, y=level, xref="paper", yref="y3",
                text=_fmt_price(level), showarrow=False, xanchor="right",
                font=dict(color="rgba(148,163,184,0.45)", size=9, family="Inter, sans-serif"),
                bgcolor="rgba(0,0,0,0)",
            ))

        # ── Row 4: Volume ───────────────────────────────────────────────

        if "Volume" in df.columns:
            vol_colors = ["rgba(38,166,154,0.5)" if c >= o else "rgba(239,83,80,0.5)" for c, o in zip(close, opn)]
            fig.add_trace(go.Bar(
                x=dates, y=df["Volume"], marker_color=vol_colors,
                marker_line_width=0,
                name="Volume", showlegend=False,
                hovertemplate="Vol: %{y:,.0f}<extra></extra>",
            ), row=4, col=1)

        # ── Row labels ──────────────────────────────────────────────────

        for label, yref in [("MACD (12,26,9)", "y"), ("Stoch (14,3,3)", "y2")]:
            annotations.append(dict(
                x=0.005, y=1.0, xref="paper", yref=f"{yref} domain",
                text=f"<b>{label}</b>", showarrow=False,
                xanchor="left", yanchor="top",
                font=dict(color="rgba(148,163,184,0.4)", size=9, family="Inter, sans-serif"),
            ))
        annotations.append(dict(
            x=0.005, y=0.0, xref="paper", yref="y4 domain",
            text="<b>Vol</b>", showarrow=False,
            xanchor="left", yanchor="bottom",
            font=dict(color="rgba(148,163,184,0.4)", size=9, family="Inter, sans-serif"),
        ))

        # ── Layout ──────────────────────────────────────────────────────

        fig.update_layout(
            annotations=annotations,
            xaxis_rangeslider_visible=False,
            hovermode="x unified",
            hoverdistance=30,
            spikedistance=-1,
            dragmode="pan",
            margin=dict(l=48, r=72, t=4, b=18),
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)",
            font=dict(family="'Inter', -apple-system, sans-serif", size=11),
            legend=dict(visible=False),
        )

        # Kill rangeslider on every x-axis (candlestick adds one by default)
        fig.update_xaxes(
            showgrid=True, gridcolor="rgba(148,163,184,0.035)",
            zeroline=False, showspikes=False,
            range=x_range_initial,
            rangeslider=dict(visible=False),
        )
        fig.update_yaxes(
            showgrid=True, gridcolor="rgba(148,163,184,0.035)",
            zeroline=False, showspikes=False,
            tickfont=dict(size=9),
            side="right",
        )
        # Only show x-tick labels on the bottom row
        for r in [1, 2, 3]:
            fig.update_xaxes(showticklabels=False, row=r, col=1)
        fig.update_xaxes(tickfont=dict(size=9), row=4, col=1)

        # Subtle separator lines between panels
        for yref in ["y", "y2", "y4"]:
            fig.add_shape(
                type="line", x0=0, x1=1, y0=0, y1=0,
                xref="paper", yref=f"{yref} domain",
                line=dict(color="rgba(148,163,184,0.1)", width=0.5),
            )

        return json.loads(pio.to_json(fig, engine="json"))

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Research chart failed for ticker=%s: %s", ticker, e)
        raise HTTPException(status_code=500, detail=f"Research chart error: {e}")

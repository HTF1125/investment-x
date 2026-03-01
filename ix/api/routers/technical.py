from __future__ import annotations

import json
import os
import textwrap
from io import BytesIO
from typing import Dict, List, Optional, Any
from datetime import datetime, timedelta

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import yfinance as yf
from fastapi import APIRouter, HTTPException, Query, Body
from fastapi.responses import StreamingResponse
from plotly.subplots import make_subplots
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import utils
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import markdown
from PIL import Image

router = APIRouter()


class TDSequentialClean:
    def __init__(
        self,
        show_setup_from: int = 9,
        show_countdown_from: int = 13,
        label_cooldown_bars: int = 0,
        setup_lookback: int = 4,
        setup_len: int = 9,
        countdown_lookback: int = 2,
        countdown_len: int = 13,
        suppress_setup_when_cd_active: bool = False,
        cancel_on_opposite_setup9: bool = True,
    ):
        self.show_setup_from = show_setup_from
        self.show_countdown_from = show_countdown_from
        self.label_cooldown_bars = label_cooldown_bars
        self.setup_lookback = setup_lookback
        self.setup_len = setup_len
        self.countdown_lookback = countdown_lookback
        self.countdown_len = countdown_len
        self.suppress_setup_when_cd_active = suppress_setup_when_cd_active
        self.cancel_on_opposite_setup9 = cancel_on_opposite_setup9

    def _cooldown_mask(self, base_mask: pd.Series) -> pd.Series:
        if self.label_cooldown_bars <= 0:
            return base_mask
        keep = np.zeros(len(base_mask), dtype=bool)
        cooldown = 0
        for i, flag in enumerate(base_mask.to_numpy()):
            if cooldown > 0:
                cooldown -= 1
                continue
            if flag:
                keep[i] = True
                cooldown = self.label_cooldown_bars
        return pd.Series(keep, index=base_mask.index)

    def compute(self, close: pd.Series) -> pd.DataFrame:
        close = pd.to_numeric(close, errors="coerce").astype(float).dropna().sort_index()
        if len(close) < 20:
            raise ValueError("close series too short")
        n = len(close)
        bear_cond = close > close.shift(self.setup_lookback)
        bull_cond = close < close.shift(self.setup_lookback)
        bear_setup = np.zeros(n, dtype=int)
        bull_setup = np.zeros(n, dtype=int)
        bs = us = 0
        for i in range(n):
            bs = bs + 1 if bool(bear_cond.iloc[i]) else 0
            us = us + 1 if bool(bull_cond.iloc[i]) else 0
            bear_setup[i] = bs if 1 <= bs <= self.setup_len else 0
            bull_setup[i] = us if 1 <= us <= self.setup_len else 0

        bear_cd_cond = close >= close.shift(self.countdown_lookback)
        bull_cd_cond = close <= close.shift(self.countdown_lookback)
        bear_cd = np.zeros(n, dtype=int)
        bull_cd = np.zeros(n, dtype=int)
        bear_cd_active = np.zeros(n, dtype=bool)
        bull_cd_active = np.zeros(n, dtype=bool)
        bear_active = bull_active = False
        bear_count = bull_count = 0
        for i in range(n):
            if bear_setup[i] == self.setup_len:
                bear_active = True
                bear_count = 0
                if self.cancel_on_opposite_setup9:
                    bull_active = False
                    bull_count = 0
            if bull_setup[i] == self.setup_len:
                bull_active = True
                bull_count = 0
                if self.cancel_on_opposite_setup9:
                    bear_active = False
                    bear_count = 0
            bear_cd_active[i] = bear_active
            bull_cd_active[i] = bull_active
            if bear_active:
                if bear_count < self.countdown_len and bool(bear_cd_cond.iloc[i]):
                    bear_count += 1
                bear_cd[i] = bear_count
                if bear_count >= self.countdown_len:
                    bear_active = False
            if bull_active:
                if bull_count < self.countdown_len and bool(bull_cd_cond.iloc[i]):
                    bull_count += 1
                bull_cd[i] = bull_count
                if bull_count >= self.countdown_len:
                    bull_active = False

        return pd.DataFrame(
            {
                "close": close,
                "bear_setup": bear_setup,
                "bull_setup": bull_setup,
                "bear_cd": bear_cd,
                "bull_cd": bull_cd,
                "bear_cd_active": bear_cd_active,
                "bull_cd_active": bull_cd_active,
            },
            index=close.index,
        )


def _find_swings(df: pd.DataFrame, window: int) -> pd.DataFrame:
    high = df["High"]
    low = df["Low"]
    if isinstance(high, pd.DataFrame):
        high = high.iloc[:, 0]
    if isinstance(low, pd.DataFrame):
        low = low.iloc[:, 0]
    is_high = high.eq(high.rolling(window * 2 + 1, center=True).max())
    is_low = low.eq(low.rolling(window * 2 + 1, center=True).min())
    points = []
    for i in range(len(df)):
        d = df.index[i]
        if bool(is_high.iloc[i]):
            points.append((d, float(high.iloc[i]), "H"))
        if bool(is_low.iloc[i]):
            points.append((d, float(low.iloc[i]), "L"))
    if not points:
        return pd.DataFrame(columns=["Date", "Price", "Type"])
    swings = pd.DataFrame(points, columns=["Date", "Price", "Type"]).sort_values("Date")
    cleaned = []
    for _, p in swings.iterrows():
        if not cleaned:
            cleaned.append(p)
            continue
        last = cleaned[-1]
        if p["Type"] == last["Type"]:
            if (p["Type"] == "H" and p["Price"] > last["Price"]) or (p["Type"] == "L" and p["Price"] < last["Price"]):
                cleaned[-1] = p
        else:
            cleaned.append(p)
    return pd.DataFrame(cleaned)


def _wave_labels(n: int) -> list[str]:
    seq = ["(1)", "(2)", "(3)", "(4)", "(5)", "(a)", "(b)", "(c)"]
    return seq[:n] if n <= len(seq) else seq + [f"({i})" for i in range(6, 6 + n - len(seq))]

def _alternates_types(types: list[str]) -> bool:
    return all(types[i] != types[i - 1] for i in range(1, len(types)))


def _score_to_target(value: float, target: float, spread: float) -> float:
    if spread <= 0 or not np.isfinite(value):
        return 0.0
    return max(0.0, 1.0 - abs(value - target) / spread)


def _evaluate_motive(prices: np.ndarray, types: list[str], bullish: bool) -> tuple[bool, float]:
    # Contiguous 6 pivots: p0..p5 define waves 1..5.
    p0, p1, p2, p3, p4, p5 = [float(x) for x in prices]
    tol = max((max(prices) - min(prices)) * 0.005, 1e-9)

    if bullish:
        if types != ["L", "H", "L", "H", "L", "H"]:
            return False, float("-inf")
        w1 = p1 - p0
        w2 = p1 - p2
        w3 = p3 - p2
        w4 = p3 - p4
        w5 = p5 - p4
        if min(w1, w2, w3, w4, w5) <= 0:
            return False, float("-inf")
        if p2 <= p0 + tol:
            return False, float("-inf")
        if p4 < p1 - tol:
            return False, float("-inf")
        if p5 <= p3 + tol * 0.3:
            return False, float("-inf")
        if w3 < min(w1, w5) + tol:
            return False, float("-inf")
    else:
        if types != ["H", "L", "H", "L", "H", "L"]:
            return False, float("-inf")
        w1 = p0 - p1
        w2 = p2 - p1
        w3 = p2 - p3
        w4 = p4 - p3
        w5 = p4 - p5
        if min(w1, w2, w3, w4, w5) <= 0:
            return False, float("-inf")
        if p2 >= p0 - tol:
            return False, float("-inf")
        if p4 > p1 + tol:
            return False, float("-inf")
        if p5 >= p3 - tol * 0.3:
            return False, float("-inf")
        if w3 < min(w1, w5) + tol:
            return False, float("-inf")

    r2 = w2 / w1
    r4 = w4 / w3
    ext3 = w3 / w1
    ext5 = w5 / w1
    if not (0.15 <= r2 <= 0.95):
        return False, float("-inf")
    if not (0.10 <= r4 <= 0.90):
        return False, float("-inf")
    if ext3 < 1.0 or ext5 < 0.20:
        return False, float("-inf")

    impulse_legs = np.array([w1, w3, w5], dtype=float)
    if impulse_legs.min() < impulse_legs.max() * 0.12:
        return False, float("-inf")

    displacement = abs(p5 - p0) / (impulse_legs.sum() + 1e-9)
    score = 2.0
    score += _score_to_target(r2, target=0.618, spread=0.50)
    score += _score_to_target(r4, target=0.382, spread=0.35)
    score += _score_to_target(ext3, target=1.618, spread=1.10)
    score += _score_to_target(ext5, target=1.000, spread=0.90)
    score += min(1.0, displacement)
    return True, score


def _valid_motive(prices: np.ndarray, types: list[str], bullish: bool) -> bool:
    valid, _ = _evaluate_motive(prices, types, bullish)
    return valid


def _find_motive_segment(piv: pd.DataFrame, max_lookback: int = 120) -> tuple[int, bool, float] | None:
    if len(piv) < 6:
        return None
    recent = piv.tail(max_lookback).reset_index(drop=True)
    offset = len(piv) - len(recent)
    pivot_scale = np.nanmedian(np.abs(np.diff(recent["Price"].to_numpy(dtype=float))))
    if not np.isfinite(pivot_scale) or pivot_scale <= 0:
        pivot_scale = 1.0
    candidates: list[tuple[int, bool, float]] = []
    for i in range(0, len(recent) - 5):
        seg = recent.iloc[i : i + 6]
        types = seg["Type"].tolist()
        if not _alternates_types(types):
            continue
        prices = seg["Price"].to_numpy(dtype=float)
        for bullish in (True, False):
            valid, base_score = _evaluate_motive(prices, types, bullish)
            if valid:
                displacement = abs(float(prices[-1] - prices[0])) / pivot_scale
                recency = (i + 5) / max(6, len(recent))
                total_score = base_score + min(3.0, displacement * 0.20) + recency * 0.60
                candidates.append((i, bullish, total_score))
    if not candidates:
        return None
    # Prefer strongest structure score, then latest start.
    candidates.sort(key=lambda x: (x[2], x[0]), reverse=True)
    start_idx, bullish, total_score = candidates[0]
    return (start_idx + offset, bullish, total_score)


def _evaluate_abc(prices: np.ndarray, types: list[str], bullish_motive: bool) -> tuple[bool, float]:
    a, b, c, d = [float(x) for x in prices]  # d is prior wave-5 pivot
    tol = max(abs(d - a) * 0.005, 1e-9)
    if bullish_motive:
        # after bullish motive ends at high d, correction is L-H-L
        if types != ["L", "H", "L"]:
            return False, float("-inf")
        if not (a < d - tol and b < d - tol and c < b - tol * 0.2):
            return False, float("-inf")
        if c > a + tol:
            return False, float("-inf")
        leg1 = d - a
        leg2 = b - a
        leg3 = b - c
    else:
        # after bearish motive ends at low d, correction is H-L-H
        if types != ["H", "L", "H"]:
            return False, float("-inf")
        if not (a > d + tol and b > d + tol and c > b + tol * 0.2):
            return False, float("-inf")
        if c < a - tol:
            return False, float("-inf")
        leg1 = a - d
        leg2 = a - b
        leg3 = c - b

    if min(leg1, leg2, leg3) <= 0:
        return False, float("-inf")
    retr_b = leg2 / leg1
    ext_c = leg3 / leg2
    if not (0.15 <= retr_b <= 0.95):
        return False, float("-inf")
    if not (0.40 <= ext_c <= 2.80):
        return False, float("-inf")

    score = 1.0
    score += _score_to_target(retr_b, target=0.618, spread=0.45)
    score += _score_to_target(ext_c, target=1.000, spread=0.90)
    if bullish_motive and c <= a - tol:
        score += 0.40
    if (not bullish_motive) and c >= a + tol:
        score += 0.40
    return True, score


def _valid_abc(prices: np.ndarray, types: list[str], bullish_motive: bool) -> bool:
    valid, _ = _evaluate_abc(prices, types, bullish_motive)
    return valid


def _best_abc_after_motive(
    piv: pd.DataFrame,
    motive_end_idx: int,
    bullish_motive: bool,
    max_forward: int = 8,
) -> tuple[pd.DataFrame | None, float]:
    if motive_end_idx + 3 > len(piv):
        return None, float("-inf")
    start_min = motive_end_idx + 1
    start_max = min(len(piv) - 3, motive_end_idx + max_forward)
    if start_min > start_max:
        return None, float("-inf")

    last_p5 = float(piv["Price"].iloc[motive_end_idx])
    best: tuple[pd.DataFrame, float] | None = None
    for start in range(start_min, start_max + 1):
        abc = piv.iloc[start : start + 3].copy()
        abc_prices = np.array(
            [float(abc["Price"].iloc[0]), float(abc["Price"].iloc[1]), float(abc["Price"].iloc[2]), last_p5]
        )
        abc_types = abc["Type"].tolist()
        valid, base_score = _evaluate_abc(abc_prices, abc_types, bullish_motive=bullish_motive)
        if not valid:
            continue
        delay_penalty = (start - start_min) * 0.08
        total = base_score - delay_penalty
        if best is None or total > best[1]:
            best = (abc, total)

    if best is None:
        return None, float("-inf")
    return best


def _extract_elliott_labels(
    piv: pd.DataFrame,
) -> tuple[pd.DataFrame | None, pd.DataFrame | None, bool | None, float]:
    found = _find_motive_segment(piv)
    if found is None:
        return None, None, None, float("-inf")
    start, bullish, motive_score = found
    motive = piv.iloc[start : start + 6].copy()
    correction, correction_score = _best_abc_after_motive(
        piv=piv,
        motive_end_idx=start + 5,
        bullish_motive=bullish,
    )
    total_score = motive_score + (correction_score if correction is not None else 0.0)
    return motive, correction, bullish, total_score


def _extract_best_elliott(
    swings_map: Dict[int, pd.DataFrame]
) -> tuple[pd.DataFrame | None, pd.DataFrame | None, bool | None, int | None]:
    degree_bias = {4: 0.05, 8: 0.20, 16: 0.15}
    best: tuple[float, int, pd.DataFrame, pd.DataFrame | None, bool] | None = None
    for length, piv in swings_map.items():
        if piv is None or len(piv) < 6:
            continue
        motive, correction, bullish, score = _extract_elliott_labels(piv)
        if motive is None or bullish is None:
            continue
        adjusted = score + degree_bias.get(length, 0.0)
        if best is None or adjusted > best[0]:
            best = (adjusted, length, motive, correction, bullish)
    if best is None:
        return None, None, None, None
    _, length, motive, correction, bullish = best
    return motive, correction, bullish, length


def _add_fib_zone(fig: go.Figure, piv: pd.DataFrame) -> None:
    if len(piv) < 2:
        return
    p0 = float(piv["Price"].iloc[-2])
    p1 = float(piv["Price"].iloc[-1])
    diff = abs(p1 - p0)
    if diff <= 0:
        return
    x0 = piv["Date"].iloc[-1]
    x1 = piv["Date"].iloc[-1] + (piv["Date"].iloc[-1] - piv["Date"].iloc[-2]) * 2
    levels = [0.5, 0.618, 0.764, 0.854]
    sign = -1 if p1 > p0 else 1
    ys = [p1 + sign * diff * lv for lv in levels]
    for i, y in enumerate(ys):
        fig.add_hline(y=y, line_width=1, line_dash="dot", line_color="rgba(96,165,250,0.55)" if i < 2 else "rgba(239,68,68,0.55)", row=1, col=1)
    fig.add_shape(
        type="rect",
        xref="x",
        yref="y",
        x0=x0,
        x1=x1,
        y0=min(ys[-2:]),
        y1=max(ys[-2:]),
        fillcolor="rgba(59,130,246,0.12)",
        line=dict(color="rgba(59,130,246,0.35)", width=1),
    )


def _normalize_yf(data: pd.DataFrame, ticker: str) -> pd.DataFrame:
    if isinstance(data.columns, pd.MultiIndex):
        if ticker in data.columns.get_level_values(-1):
            data = data.xs(ticker, axis=1, level=-1)
        else:
            data.columns = data.columns.get_level_values(0)
    return data[["Open", "High", "Low", "Close", "Volume"]].dropna()


def _compute_rsi(close: pd.Series, period: int = 14) -> pd.Series:
    delta = close.diff()
    gain = delta.clip(lower=0.0)
    loss = -delta.clip(upper=0.0)
    avg_gain = gain.ewm(alpha=1 / period, min_periods=period, adjust=False).mean()
    avg_loss = loss.ewm(alpha=1 / period, min_periods=period, adjust=False).mean()
    rs = avg_gain / avg_loss.replace(0, np.nan)
    rsi = 100.0 - (100.0 / (1.0 + rs))
    return rsi.clip(lower=0.0, upper=100.0)


def _compute_squeeze_momentum(
    df: pd.DataFrame,
    bb_length: int = 20,
    bb_mult: float = 2.0,
    kc_length: int = 20,
    kc_mult: float = 1.5,
    use_true_range: bool = True,
) -> pd.DataFrame:
    """LazyBear Squeeze Momentum Indicator."""
    close = pd.to_numeric(df["Close"], errors="coerce")
    high = pd.to_numeric(df["High"], errors="coerce")
    low = pd.to_numeric(df["Low"], errors="coerce")

    # Bollinger Bands
    basis = close.rolling(bb_length).mean()
    bb_dev = kc_mult * close.rolling(bb_length).std()
    upper_bb = basis + bb_dev
    lower_bb = basis - bb_dev

    # Keltner Channels
    ma_kc = close.rolling(kc_length).mean()
    if use_true_range:
        prev_close = close.shift(1)
        tr = pd.concat([high - low, (high - prev_close).abs(), (low - prev_close).abs()], axis=1).max(axis=1)
    else:
        tr = high - low
    range_ma = tr.rolling(kc_length).mean()
    upper_kc = ma_kc + range_ma * kc_mult
    lower_kc = ma_kc - range_ma * kc_mult

    sqz_on = (lower_bb > lower_kc) & (upper_bb < upper_kc)
    sqz_off = (lower_bb < lower_kc) & (upper_bb > upper_kc)
    no_sqz = ~sqz_on & ~sqz_off

    # Momentum value: linear regression of delta source
    highest_high = high.rolling(kc_length).max()
    lowest_low = low.rolling(kc_length).min()
    mid_hl = (highest_high + lowest_low) / 2
    sma_close = close.rolling(kc_length).mean()
    delta = close - (mid_hl + sma_close) / 2

    # Linear regression (linreg) over kc_length
    val = pd.Series(index=close.index, dtype=float)
    arr = delta.to_numpy(dtype=float)
    for i in range(kc_length - 1, len(arr)):
        y = arr[i - kc_length + 1 : i + 1]
        if np.isnan(y).any():
            val.iloc[i] = np.nan
            continue
        x = np.arange(kc_length, dtype=float)
        m, b = np.polyfit(x, y, 1)
        val.iloc[i] = m * (kc_length - 1) + b

    val_prev = val.shift(1)
    colors = pd.Series("gray", index=close.index)
    colors[val > 0] = "lime"
    colors[(val > 0) & (val < val_prev)] = "green"
    colors[val < 0] = "red"
    colors[(val < 0) & (val > val_prev)] = "maroon"

    sqz_color = pd.Series("gray", index=close.index)
    sqz_color[no_sqz] = "blue"
    sqz_color[sqz_on] = "black"

    return pd.DataFrame({
        "val": val,
        "bar_color": colors,
        "sqz_on": sqz_on,
        "sqz_off": sqz_off,
        "no_sqz": no_sqz,
        "sqz_dot_color": sqz_color,
    }, index=close.index)


def _compute_supertrend(
    df: pd.DataFrame,
    atr_period: int = 10,
    multiplier: float = 3.0,
) -> pd.DataFrame:
    """Classic Supertrend indicator."""
    close = pd.to_numeric(df["Close"], errors="coerce")
    high = pd.to_numeric(df["High"], errors="coerce")
    low = pd.to_numeric(df["Low"], errors="coerce")
    hl2 = (high + low) / 2

    # ATR
    prev_close = close.shift(1)
    tr = pd.concat([high - low, (high - prev_close).abs(), (low - prev_close).abs()], axis=1).max(axis=1)
    atr = tr.ewm(span=atr_period, adjust=False).mean()

    up = hl2 - multiplier * atr
    dn = hl2 + multiplier * atr

    up_arr = up.to_numpy(dtype=float)
    dn_arr = dn.to_numpy(dtype=float)
    close_arr = close.to_numpy(dtype=float)
    trend = np.ones(len(close_arr), dtype=int)

    for i in range(1, len(close_arr)):
        if np.isnan(up_arr[i]) or np.isnan(dn_arr[i]) or np.isnan(close_arr[i]):
            trend[i] = trend[i - 1]
            continue
        # Adjust up/dn
        up_arr[i] = max(up_arr[i], up_arr[i - 1]) if close_arr[i - 1] > up_arr[i - 1] else up_arr[i]
        dn_arr[i] = min(dn_arr[i], dn_arr[i - 1]) if close_arr[i - 1] < dn_arr[i - 1] else dn_arr[i]
        if trend[i - 1] == -1 and close_arr[i] > dn_arr[i - 1]:
            trend[i] = 1
        elif trend[i - 1] == 1 and close_arr[i] < up_arr[i - 1]:
            trend[i] = -1
        else:
            trend[i] = trend[i - 1]

    trend_s = pd.Series(trend, index=close.index)
    up_s = pd.Series(np.where(trend == 1, up_arr, np.nan), index=close.index)
    dn_s = pd.Series(np.where(trend == -1, dn_arr, np.nan), index=close.index)

    buy_signal = (trend_s == 1) & (trend_s.shift(1) == -1)
    sell_signal = (trend_s == -1) & (trend_s.shift(1) == 1)

    return pd.DataFrame({
        "trend": trend_s,
        "up": up_s,
        "dn": dn_s,
        "buy": buy_signal,
        "sell": sell_signal,
    }, index=close.index)


def _compute_moving_averages(
    df: pd.DataFrame,
    configs: List[dict],
) -> List[dict]:
    """Return list of MA trace dicts. Each config: {type, period, color}."""
    close = pd.to_numeric(df["Close"], errors="coerce")
    traces = []
    for cfg in configs:
        ma_type = cfg.get("type", "SMA")
        period = int(cfg.get("period", 20))
        color = cfg.get("color", "#94a3b8")
        if ma_type == "EMA":
            values = close.ewm(span=period, adjust=False).mean()
        elif ma_type == "WMA":
            weights = np.arange(1, period + 1, dtype=float)
            values = close.rolling(period).apply(lambda x: np.dot(x, weights) / weights.sum(), raw=True)
        else:  # SMA
            values = close.rolling(period).mean()
        dates = [d.isoformat() if hasattr(d, "isoformat") else str(d) for d in df.index]
        traces.append({
            "name": f"{ma_type} {period}",
            "x": dates,
            "y": [None if np.isnan(v) else round(float(v), 4) for v in values],
            "color": color,
            "period": period,
            "type": ma_type,
        })
    return traces


def _build_figure(
    df: pd.DataFrame,
    ticker: str,
    show_setup_from: int,
    show_countdown_from: int,
    label_cooldown_bars: int,
    visible_start: Optional[pd.Timestamp] = None,
    visible_end: Optional[pd.Timestamp] = None,
    show_macd: bool = True,
    show_rsi: bool = True,
) -> go.Figure:
    rows = 1
    if show_macd: rows += 1
    if show_rsi: rows += 1
    
    # Map each indicator to its specific row index
    macd_row = 2 if show_macd else None
    rsi_row = (2 if not show_macd else 3) if show_rsi else None

    # Calculate row heights
    if rows == 3:
        row_heights = [0.62, 0.23, 0.15]
    elif rows == 2:
        row_heights = [0.75, 0.25]
    else:
        row_heights = [1.0]

    fig = make_subplots(
        rows=rows,
        cols=1,
        shared_xaxes=True,
        vertical_spacing=0.04 if rows > 1 else 0,
        row_heights=row_heights,
    )

    ohlc_hover = [
        (
            f"Date: {d.strftime('%Y-%m-%d')}<br>"
            f"Open: {o:.2f}<br>"
            f"High: {h:.2f}<br>"
            f"Low: {l:.2f}<br>"
            f"Close: {c:.2f}"
        )
        for d, o, h, l, c in zip(df.index, df["Open"], df["High"], df["Low"], df["Close"])
    ]
    fig.add_trace(
        go.Candlestick(
            x=df.index,
            open=df["Open"],
            high=df["High"],
            low=df["Low"],
            close=df["Close"],
            name=f"{ticker.upper()} Price",
            increasing_line_color="#22c55e",
            increasing_fillcolor="rgba(34,197,94,0.55)",
            decreasing_line_color="#f43f5e",
            decreasing_fillcolor="rgba(244,63,94,0.55)",
            line_width=1.0,
            whiskerwidth=0.2,
            legendgroup="price",
            text=ohlc_hover,
            hoverinfo="text",
        ),
        row=1,
        col=1,
    )

    close = pd.to_numeric(df["Close"], errors="coerce")
    ma_layers = [
        ("MA 5", 5, "#22c55e", 1.5),
        ("MA 20", 20, "#f59e0b", 1.7),
        ("MA 200", 200, "#8b5cf6", 1.9),
    ]
    for name, window, color, width in ma_layers:
        ma = close.rolling(window=window, min_periods=window).mean()
        if ma.notna().any():
            fig.add_trace(
                go.Scatter(
                    x=df.index,
                    y=ma,
                    mode="lines",
                    line=dict(color=color, width=width),
                    name=name,
                    legendgroup="price",
                    hovertemplate=f"{name}: %{{y:.2f}}<extra></extra>",
                ),
                row=1,
                col=1,
            )

    close_values = close.to_numpy(dtype=float)
    dy = np.nanmedian(np.abs(np.diff(close_values))) if len(close_values) > 2 else np.nanstd(close_values) * 0.01
    if not np.isfinite(dy) or dy <= 0:
        dy = max(np.nanstd(close_values) * 0.01, 1.0)

    swings_map: Dict[int, pd.DataFrame] = {4: _find_swings(df, 4), 8: _find_swings(df, 8), 16: _find_swings(df, 16)}
    motive, correction, bullish, wave_degree = _extract_best_elliott(swings_map)

    if wave_degree and wave_degree in swings_map and len(swings_map[wave_degree]) >= 2:
        piv = swings_map[wave_degree].tail(36)
        fig.add_trace(
            go.Scatter(
                x=piv["Date"],
                y=piv["Price"],
                mode="lines",
                line=dict(color="rgba(148,163,184,0.35)", width=1.0, dash="dot"),
                name="Wave Backbone",
                showlegend=False,
                hovertemplate="Backbone: %{y:.2f}<extra></extra>",
            ),
            row=1,
            col=1,
        )

    if motive is not None:
        wave_color = "#22c55e" if bullish else "#f43f5e"
        fig.add_trace(
            go.Scatter(
                x=motive["Date"],
                y=motive["Price"],
                mode="lines+markers",
                line=dict(color=wave_color, width=2.1),
                marker=dict(size=6, color=wave_color),
                name="Elliott 1-5",
                legendgroup="elliott",
                hovertemplate="Wave: %{y:.2f}<extra></extra>",
            ),
            row=1,
            col=1,
        )
        motive_y = [
            float(price + (dy * 1.35 if t == "H" else -dy * 1.35))
            for price, t in zip(motive["Price"], motive["Type"])
        ]
        fig.add_trace(
            go.Scatter(
                x=motive["Date"],
                y=motive_y,
                mode="text",
                text=["0", "1", "2", "3", "4", "5"],
                textfont=dict(size=12, color=wave_color),
                showlegend=False,
                hoverinfo="skip",
            ),
            row=1,
            col=1,
        )

    if correction is not None:
        corr_color = "#38bdf8" if bullish else "#fb7185"
        fig.add_trace(
            go.Scatter(
                x=correction["Date"],
                y=correction["Price"],
                mode="lines+markers",
                line=dict(color=corr_color, width=1.9),
                marker=dict(size=5, color=corr_color),
                name="Elliott A-B-C",
                legendgroup="elliott",
                hovertemplate="Correction: %{y:.2f}<extra></extra>",
            ),
            row=1,
            col=1,
        )
        corr_y = [
            float(price + (dy * 1.1 if t == "H" else -dy * 1.1))
            for price, t in zip(correction["Price"], correction["Type"])
        ]
        fig.add_trace(
            go.Scatter(
                x=correction["Date"],
                y=corr_y,
                mode="text",
                text=["A", "B", "C"],
                textfont=dict(size=11, color=corr_color),
                showlegend=False,
                hoverinfo="skip",
            ),
            row=1,
            col=1,
        )

    td = TDSequentialClean(
        show_setup_from=show_setup_from,
        show_countdown_from=show_countdown_from,
        label_cooldown_bars=label_cooldown_bars,
    )
    tdf = td.compute(df["Close"])
    above = tdf["close"] + dy * 1.9
    below = tdf["close"] - dy * 1.9
    bear_setup_mask = td._cooldown_mask(tdf["bear_setup"] >= td.show_setup_from)
    bull_setup_mask = td._cooldown_mask(tdf["bull_setup"] >= td.show_setup_from)
    bear_cd_mask = td._cooldown_mask(tdf["bear_cd"] >= td.show_countdown_from)
    bull_cd_mask = td._cooldown_mask(tdf["bull_cd"] >= td.show_countdown_from)
    if bear_setup_mask.any():
        fig.add_trace(
            go.Scatter(
                x=tdf.index[bear_setup_mask],
                y=above[bear_setup_mask],
                mode="text",
                text=tdf.loc[bear_setup_mask, "bear_setup"].astype(int).astype(str),
                textfont=dict(color="#fb7185", size=10),
                showlegend=False,
                hoverinfo="skip",
            ),
            row=1,
            col=1,
        )
    if bull_setup_mask.any():
        fig.add_trace(
            go.Scatter(
                x=tdf.index[bull_setup_mask],
                y=below[bull_setup_mask],
                mode="text",
                text=tdf.loc[bull_setup_mask, "bull_setup"].astype(int).astype(str),
                textfont=dict(color="#34d399", size=10),
                showlegend=False,
                hoverinfo="skip",
            ),
            row=1,
            col=1,
        )
    if bear_cd_mask.any():
        fig.add_trace(
            go.Scatter(
                x=tdf.index[bear_cd_mask],
                y=above[bear_cd_mask] + dy,
                mode="text",
                text=tdf.loc[bear_cd_mask, "bear_cd"].astype(int).astype(str),
                textfont=dict(color="#f43f5e", size=11),
                showlegend=False,
                hoverinfo="skip",
            ),
            row=1,
            col=1,
        )
    if bull_cd_mask.any():
        fig.add_trace(
            go.Scatter(
                x=tdf.index[bull_cd_mask],
                y=below[bull_cd_mask] - dy,
                mode="text",
                text=tdf.loc[bull_cd_mask, "bull_cd"].astype(int).astype(str),
                textfont=dict(color="#10b981", size=11),
                showlegend=False,
                hoverinfo="skip",
            ),
            row=1,
            col=1,
        )

    # MACD Subplot
    if show_macd:
        ema12 = close.ewm(span=12, adjust=False).mean()
        ema26 = close.ewm(span=26, adjust=False).mean()
        macd_line = ema12 - ema26
        macd_signal = macd_line.ewm(span=9, adjust=False).mean()
        macd_hist = macd_line - macd_signal
        hist_prev = macd_hist.shift(1)
        hist_colors = np.where(
            macd_hist >= 0,
            np.where(macd_hist >= hist_prev, "rgba(16,185,129,0.85)", "rgba(16,185,129,0.45)"),
            np.where(macd_hist <= hist_prev, "rgba(244,63,94,0.85)", "rgba(244,63,94,0.45)"),
        )
        fig.add_trace(
            go.Bar(
                x=df.index,
                y=macd_hist,
                marker_color=hist_colors,
                name="MACD Hist",
                showlegend=False,
                hovertemplate="MACD Hist: %{y:.3f}<extra></extra>",
            ),
            row=macd_row,
            col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=df.index,
                y=macd_line,
                mode="lines",
                line=dict(color="#38bdf8", width=1.8),
                name="MACD",
                legendgroup="macd",
                hovertemplate="MACD: %{y:.3f}<extra></extra>",
            ),
            row=macd_row,
            col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=df.index,
                y=macd_signal,
                mode="lines",
                line=dict(color="#f59e0b", width=1.5),
                name="Signal",
                legendgroup="macd",
                hovertemplate="Signal: %{y:.3f}<extra></extra>",
            ),
            row=macd_row,
            col=1,
        )
        fig.add_hline(y=0, line_width=1, line_color="rgba(148,163,184,0.55)", row=macd_row, col=1)
        fig.update_yaxes(title_text="MACD", row=macd_row, col=1)

    # RSI Subplot
    if show_rsi:
        rsi = _compute_rsi(close, period=14)
        rsi_mean = rsi.rolling(window=9, min_periods=9).mean()
        fig.add_trace(
            go.Scatter(
                x=df.index,
                y=rsi,
                mode="lines",
                line=dict(color="#60a5fa", width=1.8),
                name="RSI 14",
                legendgroup="rsi",
                hovertemplate="RSI 14: %{y:.1f}<extra></extra>",
            ),
            row=rsi_row,
            col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=df.index,
                y=rsi_mean,
                mode="lines",
                line=dict(color="#f59e0b", width=1.4, dash="dot"),
                name="RSI Mean",
                legendgroup="rsi",
                hovertemplate="RSI Mean: %{y:.1f}<extra></extra>",
            ),
            row=rsi_row,
            col=1,
        )
        fig.add_hline(y=70, line_width=1, line_color="rgba(244,63,94,0.75)", line_dash="dot", row=rsi_row, col=1)
        fig.add_hline(y=50, line_width=1, line_color="rgba(148,163,184,0.45)", line_dash="dash", row=rsi_row, col=1)
        fig.add_hline(y=30, line_width=1, line_color="rgba(34,197,94,0.75)", line_dash="dot", row=rsi_row, col=1)
        fig.update_yaxes(
            title_text="RSI",
            range=[0, 100],
            tickmode="array",
            tickvals=[30, 50, 70],
            ticktext=["30", "50", "70"],
            row=rsi_row,
            col=1,
        )

    # Visible window controls display range only (does not slice source data).
    if visible_start is not None and visible_end is not None and visible_start < visible_end:
        mask = (df.index >= visible_start) & (df.index <= visible_end)
        if mask.any():
            sub = df.loc[mask]
            y1_min = float(sub["Low"].min())
            y1_max = float(sub["High"].max())
            if np.isfinite(y1_min) and np.isfinite(y1_max) and y1_max > y1_min:
                pad1 = (y1_max - y1_min) * 0.08
                fig.update_yaxes(range=[y1_min - pad1, y1_max + pad1], row=1, col=1)

            if show_macd:
                ema12 = close.ewm(span=12, adjust=False).mean()
                ema26 = close.ewm(span=26, adjust=False).mean()
                macd_line = ema12 - ema26
                macd_signal = macd_line.ewm(span=9, adjust=False).mean()
                macd_hist = macd_line - macd_signal
                macd_sub = pd.concat(
                    [macd_line.loc[mask], macd_signal.loc[mask], macd_hist.loc[mask]],
                    axis=1,
                ).dropna(how="all")
                if not macd_sub.empty:
                    y2_min = float(macd_sub.min().min())
                    y2_max = float(macd_sub.max().max())
                    if np.isfinite(y2_min) and np.isfinite(y2_max):
                        if y2_max == y2_min:
                            y2_min -= 1.0
                            y2_max += 1.0
                        pad2 = (y2_max - y2_min) * 0.18
                        fig.update_yaxes(range=[y2_min - pad2, y2_max + pad2], row=macd_row, col=1)

            for r in range(1, rows + 1):
                fig.update_xaxes(range=[visible_start, visible_end], row=r, col=1)

    fig.update_layout(
        template=None,
        paper_bgcolor="#050913",
        plot_bgcolor="#070d1a",
        font=dict(family="Ubuntu, Inter, Roboto, sans-serif", color="#dbeafe", size=12),
        margin=dict(l=72, r=22, t=60, b=24),
        title=dict(text=""),
        hovermode="x unified",
        hoverlabel=dict(
            bgcolor="rgba(15,23,42,0.92)",
            bordercolor="rgba(148,163,184,0.35)",
            font=dict(color="#e2e8f0", size=11),
        ),
        legend=dict(
            orientation="h",
            x=0.01,
            y=1.02,
            xanchor="left",
            yanchor="bottom",
            bgcolor="rgba(15,23,42,0.68)",
            bordercolor="rgba(148,163,184,0.35)",
            borderwidth=1,
            font=dict(size=10),
            traceorder="normal",
            itemclick="toggleothers",
            itemdoubleclick="toggle",
        ),
        xaxis_rangeslider_visible=False,
        bargap=0.12,
    )
    fig.update_xaxes(
        showgrid=True,
        gridcolor="rgba(148,163,184,0.12)",
        linecolor="rgba(226,232,240,0.65)",
        showline=True,
        mirror=True,
        rangebreaks=[dict(bounds=["sat", "mon"])],
        showspikes=False,
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor="rgba(148,163,184,0.12)",
        linecolor="rgba(226,232,240,0.65)",
        showline=True,
        mirror=True,
        automargin=True,
        showspikes=False,
    )
    fig.update_yaxes(title_text="Price", row=1, col=1)
    return fig


@router.get("/technical/elliott")
def technical_elliott(
    ticker: str = Query("SPY"),
    period: str = Query("2y"),
    interval: str = Query("1d"),
    start: Optional[str] = Query(None),
    end: Optional[str] = Query(None),
    setup_from: int = Query(9, ge=1, le=9),
    countdown_from: int = Query(13, ge=1, le=13),
    label_cooldown: int = Query(0, ge=0, le=20),
    show_macd: bool = Query(True),
    show_rsi: bool = Query(True),
):
    try:
        tk = ticker.strip().upper()
        raw = yf.download(tk, period=period, interval=interval, auto_adjust=False, progress=False)
        if raw is None or raw.empty:
            raise HTTPException(status_code=404, detail=f"No data for ticker '{ticker}'.")
        df = _normalize_yf(raw, tk)
        if df.empty:
            raise HTTPException(status_code=404, detail=f"No OHLC data for ticker '{ticker}'.")
        vis_start = pd.to_datetime(start) if start else None
        vis_end = pd.to_datetime(end) if end else None
        fig = _build_figure(
            df,
            tk,
            setup_from,
            countdown_from,
            label_cooldown,
            visible_start=vis_start,
            visible_end=vis_end,
            show_macd=show_macd,
            show_rsi=show_rsi,
        )
        return json.loads(pio.to_json(fig, engine="json"))
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to build technical chart: {e}")


@router.get("/technical/summary")
def get_technical_summary(
    ticker: str = Query(...),
    interval: str = Query("1d"),
):
    from ix.misc.openai import TechnicalAnalyzer
    import os
    
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        return {"summary": "AI Technical Analysis is unavailable: OPENAI_API_KEY not configured."}

    try:
        tk = ticker.strip().upper()
        # Get latest 60 bars for context
        df = yf.download(tk, period="1y", interval=interval, auto_adjust=False, progress=False)
        if df.empty:
            return {"summary": f"No data available for {tk} to analyze."}
        
        df = _normalize_yf(df, tk)
        latest = df.tail(60)
        
        # Calculate some basic signals for the prompt
        ma20 = latest["Close"].rolling(20).mean().iloc[-1]
        ma50 = latest["Close"].rolling(50).mean().iloc[-1]
        rsi = _compute_rsi(latest["Close"]).iloc[-1]
        
        data_summary = f"""
        Latest Price: {latest['Close'].iloc[-1]:.2f}
        MA20: {ma20:.2f}
        MA50: {ma50:.2f}
        RSI(14): {rsi:.1f}
        Last 5 Days Close: {latest['Close'].tail(5).tolist()}
        """
        
        analyzer = TechnicalAnalyzer(api_key=api_key)
        summary_md = analyzer.analyze(tk, data_summary)
        
        return {"summary": summary_md}
    except Exception as e:
        return {"summary": f"Failed to generate AI summary: {str(e)}"}


@router.get("/technical/overlays")
def technical_overlays(
    ticker: str = Query("SPY"),
    interval: str = Query("1d"),
    # Squeeze Momentum params
    sqz: bool = Query(False),
    sqz_bb_len: int = Query(20, ge=5, le=50),
    sqz_bb_mult: float = Query(2.0, ge=0.5, le=5.0),
    sqz_kc_len: int = Query(20, ge=5, le=50),
    sqz_kc_mult: float = Query(1.5, ge=0.5, le=5.0),
    # Supertrend params
    st: bool = Query(False),
    st_period: int = Query(10, ge=2, le=50),
    st_mult: float = Query(3.0, ge=0.5, le=10.0),
    # Moving averages — comma-separated list of "TYPE:period:color"
    # e.g. "SMA:20:#f59e0b,EMA:50:#38bdf8"
    mas: Optional[str] = Query(None),
):
    """Return overlay data (Squeeze Momentum, Supertrend, extra MAs) as raw series."""
    try:
        tk = ticker.strip().upper()
        period_map = {"1d": "2y", "1wk": "5y", "1mo": "10y"}
        yf_period = period_map.get(interval, "2y")
        raw = yf.download(tk, period=yf_period, interval=interval, auto_adjust=False, progress=False)
        if raw is None or raw.empty:
            raise HTTPException(status_code=404, detail=f"No data for ticker '{tk}'.")
        df = _normalize_yf(raw, tk)
        if df.empty:
            raise HTTPException(status_code=404, detail=f"No OHLC data for ticker '{tk}'.")

        dates = [d.isoformat() if hasattr(d, "isoformat") else str(d) for d in df.index]
        result: dict = {"dates": dates}

        if sqz:
            sqz_df = _compute_squeeze_momentum(
                df,
                bb_length=sqz_bb_len,
                bb_mult=sqz_bb_mult,
                kc_length=sqz_kc_len,
                kc_mult=sqz_kc_mult,
            )
            result["squeeze"] = {
                "val": [None if np.isnan(v) else round(float(v), 6) for v in sqz_df["val"]],
                "bar_color": sqz_df["bar_color"].tolist(),
                "sqz_on": sqz_df["sqz_on"].tolist(),
                "sqz_off": sqz_df["sqz_off"].tolist(),
                "no_sqz": sqz_df["no_sqz"].tolist(),
                "sqz_dot_color": sqz_df["sqz_dot_color"].tolist(),
            }

        if st:
            st_df = _compute_supertrend(df, atr_period=st_period, multiplier=st_mult)
            result["supertrend"] = {
                "trend": st_df["trend"].tolist(),
                "up": [None if np.isnan(v) else round(float(v), 4) for v in st_df["up"]],
                "dn": [None if np.isnan(v) else round(float(v), 4) for v in st_df["dn"]],
                "buy": st_df["buy"].tolist(),
                "sell": st_df["sell"].tolist(),
            }

        if mas:
            ma_configs = []
            for part in mas.split(","):
                parts = part.strip().split(":")
                if len(parts) >= 2:
                    ma_configs.append({
                        "type": parts[0].upper(),
                        "period": int(parts[1]),
                        "color": parts[2] if len(parts) > 2 else "#94a3b8",
                    })
            if ma_configs:
                result["moving_averages"] = _compute_moving_averages(df, ma_configs)

        return result
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to compute overlays: {e}")


def _render_chart_to_image(fig: go.Figure) -> BytesIO:
    """Render Plotly figure to a high-res PNG buffer."""
    img_bytes = pio.to_image(fig, format="png", width=1200, height=700, scale=2)
    return BytesIO(img_bytes)


def _clean_markdown(md_text: str) -> str:
    """Simple cleanup of markdown symbols for basic PDF/PPTX rendering."""
    # Remove bold/italic markers
    text = md_text.replace("**", "").replace("*", "").replace("__", "").replace("_", "")
    # Remove header markers but keep text
    lines = text.split("\n")
    cleaned_lines = []
    for line in lines:
        cleaned_lines.append(line.lstrip("#").strip())
    return "\n".join(cleaned_lines)


@router.post("/technical/export")
def export_technical_report(
    ticker: str = Query(...),
    format: str = Query("pdf"), # "pdf" or "pptx"
    summary: str = Body(..., embed=True),
    # Chart params to reproduce the image
    interval: str = Query("1d"),
    setup_from: int = Query(9),
    countdown_from: int = Query(13),
    label_cooldown: int = Query(0),
    show_macd: bool = Query(True),
    show_rsi: bool = Query(True),
):
    try:
        tk = ticker.strip().upper()
        # 1. Generate the chart image
        raw = yf.download(tk, period="2y", interval=interval, auto_adjust=False, progress=False)
        if raw is None or raw.empty:
            raise HTTPException(status_code=404, detail=f"No data for ticker '{ticker}'.")
        df = _normalize_yf(raw, tk)
        
        fig = _build_figure(
            df, tk, setup_from, countdown_from, label_cooldown,
            show_macd=show_macd, show_rsi=show_rsi
        )
        # Apply export-friendly theme
        fig.update_layout(
            template="plotly_white",
            paper_bgcolor="white",
            plot_bgcolor="#fcfcfc",
            font=dict(color="black", size=14),
            margin=dict(l=80, r=40, t=100, b=80),
        )
        chart_img = _render_chart_to_image(fig)

        filename = f"InvestmentX_{tk}_Analysis_{datetime.now().strftime('%Y%m%d')}.{format}"

        if format.lower() == "pptx":
            prs = Presentation()
            
            # Slide 1: Title
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"Technical Analysis: {tk}"
            subtitle.text = f"Generated by Investment-X Engine\nDate: {datetime.now().strftime('%Y-%m-%d')}"

            # Slide 2: Chart
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(chart_img, Inches(0.5), Inches(1), width=Inches(9))
            
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = f"{tk} Price Action & Indicators ({interval})"
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(24)

            # Slide 3+: Intelligence Report (Split if too long)
            content_lines = _clean_markdown(summary).split("\n")
            chunks = [content_lines[i:i + 15] for i in range(0, len(content_lines), 15)]
            
            for i, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Intelligence Report {' (cont.)' if i > 0 else ''}"
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                tf.text = "\n".join(chunk)
                for p in tf.paragraphs:
                    p.font.size = Pt(14)

            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )

        else: # Default to PDF
            buffer = BytesIO()
            c = canvas.Canvas(buffer, pagesize=letter)
            width, height = letter

            # Header
            c.setFont("Helvetica-Bold", 24)
            c.drawString(50, height - 60, f"Technical Analysis: {tk}")
            c.setFont("Helvetica", 12)
            c.setStrokeColorRGB(0.2, 0.5, 0.8)
            c.line(50, height - 75, width - 50, height - 75)
            
            c.drawString(50, height - 95, f"Ticker: {tk} | Interval: {interval} | Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

            # Chart Image (Centered)
            img_reader = utils.ImageReader(chart_img)
            c.drawImage(img_reader, 50, height - 420, width=width-100, height=300, preserveAspectRatio=True)

            # Report Text
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 450, "Intelligence Report")
            c.line(50, height - 455, 200, height - 455)

            c.setFont("Helvetica", 10)
            text_object = c.beginText(50, height - 480)
            text_object.setLeading(14)
            
            wrapped_text = ""
            for line in _clean_markdown(summary).split("\n"):
                wrapped_text += "\n".join(textwrap.wrap(line, width=95)) + "\n"
            
            for line in wrapped_text.split("\n"):
                if text_object.getY() < 50:
                    c.drawText(text_object)
                    c.showPage()
                    text_object = c.beginText(50, height - 50)
                    text_object.setFont("Helvetica", 10)
                    text_object.setLeading(14)
                text_object.textLine(line)
            
            c.drawText(text_object)
            c.showPage()
            c.save()
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )

    except Exception as e:
        import traceback
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

"""Chart pack report generation (PDF, PPTX, slides)."""

import base64
import concurrent.futures
import html as html_mod
from datetime import datetime
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

from fastapi import HTTPException

from ix.common import get_logger

try:
    from xhtml2pdf import pisa
except ImportError:
    pisa = None

logger = get_logger(__name__)

_DEFAULT_DISCLAIMER = (
    "This material is provided for informational purposes only and does not "
    "constitute investment advice, a recommendation, or an offer to buy or sell "
    "any securities. The views expressed herein are those of the author(s) as of "
    "the date of publication and are subject to change without notice.\n\n"
    "Past performance is not indicative of future results. All investments involve "
    "risk, including the possible loss of principal. No representation or warranty "
    "is made regarding the accuracy or completeness of the information contained "
    "herein.\n\n"
    "This document is confidential and intended solely for the use of the "
    "addressee(s). Unauthorized distribution, reproduction, or use of this "
    "material is strictly prohibited."
)


def _require_pdf_dep() -> None:
    if pisa is None:
        raise HTTPException(
            status_code=503,
            detail="PDF export dependency is unavailable. Install `xhtml2pdf`.",
        )


def _resolve_chart_figure(
    chart: dict,
    charts_by_id: Dict[str, Any],
) -> Optional[Dict[str, Any]]:
    """Resolve a pack chart config dict to a Plotly figure dict, or None."""
    if chart.get("deleted"):
        return None
    # 1. Inline cached figure
    fig = chart.get("figure")
    if fig and isinstance(fig, dict):
        return fig
    # 2. chart_id reference
    chart_id = chart.get("chart_id")
    if chart_id and chart_id in charts_by_id:
        return charts_by_id[chart_id]
    # 3. Code execution
    code = chart.get("code")
    if code and isinstance(code, str) and code.strip():
        try:
            from ix.api.routers.charts.code_execution import execute_custom_code
            return execute_custom_code(code)
        except Exception as exc:
            logger.warning("Pack report: code execution failed: %s", exc)
            return None
    return None


def generate_pack_report_pdf(
    pack_sections: List[Tuple[str, str, List[Dict[str, Any]]]],
    theme: str = "light",
) -> BytesIO:
    """Generate a PDF from resolved pack sections.

    Args:
        pack_sections: list of (pack_name, pack_description, [{"title", "description", "figure"}])
        theme: "light" or "dark"
    """
    from ix.api.routers.charts.code_execution import render_chart_image

    _require_pdf_dep()
    buffer = BytesIO()

    # Flatten all figures for parallel rendering
    all_figures: List[Tuple[int, int, Dict[str, Any]]] = []  # (section_idx, chart_idx, figure)
    for si, (_, _, charts) in enumerate(pack_sections):
        for ci, c in enumerate(charts):
            if c.get("figure"):
                all_figures.append((si, ci, c["figure"]))

    if not all_figures:
        return buffer

    # Render images in parallel
    rendered: Dict[Tuple[int, int], Optional[bytes]] = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
        future_map = {
            executor.submit(render_chart_image, fig, theme): (si, ci)
            for si, ci, fig in all_figures
        }
        try:
            for future in concurrent.futures.as_completed(future_map, timeout=600):
                key = future_map[future]
                try:
                    rendered[key] = future.result()
                except Exception as exc:
                    logger.error("Pack report render failed %s: %s", key, exc)
                    rendered[key] = None
        except concurrent.futures.TimeoutError:
            logger.error("Pack report PDF generation timed out")

    # Theme colors
    is_dark = theme.lower() == "dark"
    body_bg = "#0f172a" if is_dark else "#ffffff"
    text_color = "#f1f5f9" if is_dark else "#1e293b"
    h1_color = "#f8fafc" if is_dark else "#0f172a"
    meta_color = "#94a3b8" if is_dark else "#64748b"
    desc_color = "#cbd5e1" if is_dark else "#475569"
    nav_accent = "#818cf8" if is_dark else "#4f46e5"
    nav_bg = "#1e293b" if is_dark else "#f1f5f9"
    nav_border = "#334155" if is_dark else "#e2e8f0"
    section_bg = "#1e293b" if is_dark else "#f8fafc"
    section_border = "#334155" if is_dark else "#e2e8f0"

    total_charts = sum(len(charts) for _, _, charts in pack_sections if charts)

    parts = [f"""<!DOCTYPE html>
<html><head><meta charset="utf-8" />
<style>
    @page {{ size: landscape; margin: 1cm 1.5cm; }}
    body {{ font-family: 'Courier New', Courier, monospace; background-color: {body_bg}; color: {text_color}; line-height: 1.4; font-size: 10pt; }}
    h2 {{ color: {h1_color}; font-size: 14pt; margin-top: 6px; margin-bottom: 3px; }}
    .meta {{ color: {meta_color}; font-size: 8pt; margin-bottom: 3px; }}
    .desc {{ color: {desc_color}; font-size: 9pt; margin-bottom: 6px; font-style: italic; }}
    img {{ max-width: 100%; display: block; margin: 4px auto; }}
    .chart-page {{ page-break-before: always; }}
    .cover {{ text-align: center; padding-top: 160px; }}
    .cover-title {{ font-size: 36pt; font-weight: bold; color: {h1_color}; margin-bottom: 20px; }}
    .section-page {{ page-break-before: always; text-align: center; padding-top: 200px; }}
    .section-title {{ font-size: 28pt; font-weight: bold; color: {h1_color}; margin-bottom: 10px; }}
    .section-desc {{ font-size: 11pt; color: {desc_color}; margin-top: 8px; font-style: italic; }}
</style></head><body>
    <div class="cover">
        <div class="cover-title">Investment-X ChartPack Report</div>
        <br/>
        <span style="color: {meta_color}; font-size: 14pt;">Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}</span>
        <br/><br/>
        <span style="color: {meta_color}; font-size: 12pt;">{len(pack_sections)} pack{'s' if len(pack_sections) != 1 else ''} &middot; {total_charts} charts</span>
    </div>"""]

    chart_num = 0
    for si, (pack_name, pack_desc, charts) in enumerate(pack_sections):
        # Section divider page
        esc_name = html_mod.escape(pack_name)
        parts.append(f'<div class="section-page">')
        parts.append(f'<div class="section-title">{esc_name}</div>')
        parts.append(f'<span style="color: {meta_color}; font-size: 12pt;">{len(charts)} chart{"s" if len(charts) != 1 else ""}</span>')
        if pack_desc:
            parts.append(f'<div class="section-desc">{html_mod.escape(pack_desc)}</div>')
        parts.append("</div>")

        for ci, c in enumerate(charts):
            chart_num += 1
            title = html_mod.escape(c.get("title") or f"Chart {ci + 1}")
            desc = html_mod.escape(c.get("description") or "").replace("\n", "<br/>")

            # Nav bar
            nav_bar = (
                f'<table width="100%" cellpadding="0" cellspacing="0" '
                f'style="background-color:{nav_bg};border-bottom:1px solid {nav_border};'
                f'padding:6px 10px;margin-bottom:6px;">'
                f"<tr>"
                f'<td style="padding:4px 8px;">'
                f'<span style="font-size:10pt;font-weight:bold;color:{nav_accent};letter-spacing:1px;">'
                f"{esc_name}</span></td>"
                f'<td style="text-align:right;padding:4px 8px;">'
                f'<span style="color:{meta_color};font-size:9pt;">'
                f"{chart_num} / {total_charts}</span></td>"
                f"</tr></table>"
            )

            parts.append(f'<div class="chart-page">')
            parts.append(nav_bar)
            parts.append(f"<h2>{title}</h2>")
            if desc:
                parts.append(f'<div class="desc">{desc}</div>')

            img_bytes = rendered.get((si, ci))
            if img_bytes:
                b64 = base64.b64encode(img_bytes).decode("utf-8")
                parts.append(f'<img src="data:image/png;base64,{b64}" />')
            else:
                parts.append(
                    f'<div class="meta" style="color: red;">[Chart rendering failed]</div>'
                )
            parts.append("</div>")

    parts.append("</body></html>")

    pisa_status = pisa.CreatePDF("".join(parts), dest=buffer)
    if pisa_status.err:
        logger.error("pisa.CreatePDF failed for pack report")

    buffer.seek(0)
    return buffer


def generate_pack_report_pptx(
    slides: List[Dict[str, Any]],
    theme: str = "light",
    report_title: str = "Investment-X Report",
    classification: str = "For Internal Use Only",
) -> BytesIO:
    """Generate a layout-aware PPTX from slide data.

    Each slide dict has: layout, title, subtitle, narrative, figure, figure2, figure3,
    kpis, agendaItems, columns.
    """
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from ix.api.routers.charts.code_execution import render_chart_image

    buffer = BytesIO()
    prs = PptxPresentation()

    # 16:9 landscape
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    is_dark = theme.lower() == "dark"
    C = {  # color palette
        "bg": RGBColor(0x0F, 0x17, 0x2A) if is_dark else RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0xF8, 0xFA, 0xFC) if is_dark else RGBColor(0x0F, 0x17, 0x2A),
        "text": RGBColor(0xE2, 0xE8, 0xF0) if is_dark else RGBColor(0x1E, 0x29, 0x3B),
        "meta": RGBColor(0x94, 0xA3, 0xB8) if is_dark else RGBColor(0x64, 0x74, 0x8B),
        "accent": RGBColor(0x63, 0x82, 0xFF) if is_dark else RGBColor(0x32, 0x50, 0xD2),
        "kpi_up": RGBColor(0x34, 0xD3, 0x99),
        "kpi_down": RGBColor(0xF8, 0x71, 0x71),
    }

    def _bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = C["bg"]

    def _txt(shape, text, size=22, bold=False, color="title", align=PP_ALIGN.LEFT):
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = C[color]
        p.alignment = align
        return tf

    def _multiline(shape, text, size=11, color="text"):
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = C[color]
            p.space_after = Pt(4)
        return tf

    def _add_img(slide, img_bytes, left, top, w, h):
        if img_bytes:
            slide.shapes.add_picture(BytesIO(img_bytes), Inches(left), Inches(top), width=Inches(w), height=Inches(h))

    def _badge(slide, vi, total):
        box = slide.shapes.add_textbox(Inches(12.0), Inches(0.2), Inches(1.0), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = f"{vi + 1} / {total}"
        p.font.size = Pt(10)
        p.font.color.rgb = C["meta"]
        p.alignment = PP_ALIGN.RIGHT

    def _classification_footer(slide):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(6), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = classification
        p.font.size = Pt(7)
        p.font.color.rgb = C["meta"]

    # ── Collect all figures for parallel rendering ──
    chart_images: Dict[Tuple[int, int], Optional[bytes]] = {}
    figures_to_render: List[Tuple[int, int, dict]] = []
    for i, s in enumerate(slides):
        for slot, key in enumerate(["figure", "figure2", "figure3"]):
            fig = s.get(key)
            if fig and isinstance(fig, dict):
                figures_to_render.append((i, slot, fig))

    if figures_to_render:
        with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
            future_map = {
                executor.submit(render_chart_image, fig, theme): (idx, slot)
                for idx, slot, fig in figures_to_render
            }
            try:
                for future in concurrent.futures.as_completed(future_map, timeout=600):
                    key = future_map[future]
                    try:
                        chart_images[key] = future.result()
                    except Exception as exc:
                        logger.error("PPTX chart render failed %s: %s", key, exc)
                        chart_images[key] = None
            except concurrent.futures.TimeoutError:
                logger.error("PPTX chart rendering timed out")

    def _img(i, slot=0):
        return chart_images.get((i, slot))

    # ── Cover slide ──
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(cover)
    _txt(cover.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.2)),
         report_title, size=44, bold=True, align=PP_ALIGN.CENTER)
    tf2 = cover.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.3), Inches(0.8))
    _txt(tf2, f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}", size=18, color="meta", align=PP_ALIGN.CENTER)
    p3 = tf2.text_frame.add_paragraph()
    p3.text = f"{len(slides)} slides"
    p3.font.size = Pt(16)
    p3.font.color.rgb = C["meta"]
    p3.alignment = PP_ALIGN.CENTER

    # ── Content slides ──
    total = len(slides)
    for vi, sd in enumerate(slides):
        layout = sd.get("layout", "chart_full")
        title = sd.get("title") or ""
        subtitle_text = sd.get("subtitle") or ""
        narrative = sd.get("narrative") or ""
        kpis = sd.get("kpis") or []
        agenda = sd.get("agendaItems") or sd.get("agenda_items") or []
        columns = sd.get("columns") or []

        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(sl)
        _badge(sl, vi, total)
        _classification_footer(sl)

        if layout == "title":
            _txt(sl.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.2)),
                 title, size=44, bold=True, align=PP_ALIGN.CENTER)
            if subtitle_text:
                _txt(sl.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.3), Inches(0.6)),
                     subtitle_text, size=18, color="meta", align=PP_ALIGN.CENTER)

        elif layout == "section":
            # Accent bar
            from pptx.util import Emu
            bar = sl.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(1.2), Inches(0.08))  # MSO_SHAPE.RECTANGLE
            bar.fill.solid()
            bar.fill.fore_color.rgb = C["accent"]
            bar.line.fill.background()
            _txt(sl.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1.5)),
                 title, size=48, bold=True)
            if subtitle_text:
                _txt(sl.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(0.6)),
                     subtitle_text, size=18, color="meta")

        elif layout == "chart_full":
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=22, bold=True)
            _add_img(sl, _img(vi, 0), 0.3, 0.9, 12.7, 6.3)

        elif layout == "chart_text":
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=20, bold=True)
            _add_img(sl, _img(vi, 0), 0.3, 0.9, 8.2, 5.8)
            if narrative:
                lbl = sl.shapes.add_textbox(Inches(8.7), Inches(0.9), Inches(4.3), Inches(0.3))
                _txt(lbl, "COMMENTARY", size=8, bold=True, color="accent")
                _multiline(sl.shapes.add_textbox(Inches(8.7), Inches(1.3), Inches(4.3), Inches(5.5)),
                           narrative, size=11)

        elif layout == "text_chart":
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=20, bold=True)
            if narrative:
                _multiline(sl.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(4.3), Inches(5.8)),
                           narrative, size=11)
            _add_img(sl, _img(vi, 0), 5.0, 0.9, 8.0, 5.8)

        elif layout == "text_full":
            _txt(sl.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.3), Inches(0.8)),
                 title, size=28, bold=True)
            if narrative:
                _multiline(sl.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(4.8)),
                           narrative, size=14)

        elif layout == "two_charts":
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=22, bold=True)
            _add_img(sl, _img(vi, 0), 0.3, 0.9, 6.2, 6.0)
            _add_img(sl, _img(vi, 1), 6.8, 0.9, 6.2, 6.0)

        elif layout == "three_charts":
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=22, bold=True)
            _add_img(sl, _img(vi, 0), 0.2, 0.9, 4.2, 6.0)
            _add_img(sl, _img(vi, 1), 4.6, 0.9, 4.2, 6.0)
            _add_img(sl, _img(vi, 2), 9.0, 0.9, 4.2, 6.0)

        elif layout == "kpi_row":
            _txt(sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(0.8)),
                 title, size=28, bold=True, align=PP_ALIGN.CENTER)
            if kpis:
                n = len(kpis)
                card_w = min(2.5, 10.0 / max(n, 1))
                total_w = card_w * n + 0.3 * (n - 1)
                start_x = (13.333 - total_w) / 2
                for ki, kpi in enumerate(kpis):
                    x = start_x + ki * (card_w + 0.3)
                    # Label
                    _txt(sl.shapes.add_textbox(Inches(x), Inches(3.0), Inches(card_w), Inches(0.3)),
                         (kpi.get("label") or "").upper(), size=8, color="meta", align=PP_ALIGN.CENTER)
                    # Value
                    _txt(sl.shapes.add_textbox(Inches(x), Inches(3.4), Inches(card_w), Inches(0.7)),
                         kpi.get("value") or "—", size=28, bold=True, align=PP_ALIGN.CENTER)
                    # Change
                    change = kpi.get("change")
                    if change:
                        direction = kpi.get("direction", "flat")
                        color_key = "meta"
                        box = sl.shapes.add_textbox(Inches(x), Inches(4.2), Inches(card_w), Inches(0.3))
                        p = box.text_frame.paragraphs[0]
                        p.text = change
                        p.font.size = Pt(11)
                        p.font.color.rgb = (
                            C["kpi_up"] if direction == "up" else
                            C["kpi_down"] if direction == "down" else
                            C["meta"]
                        )
                        p.alignment = PP_ALIGN.CENTER

        elif layout == "chart_kpi":
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=22, bold=True)
            _add_img(sl, _img(vi, 0), 0.3, 0.9, 12.7, 4.5)
            if kpis:
                n = len(kpis)
                card_w = min(2.5, 10.0 / max(n, 1))
                total_w = card_w * n + 0.3 * (n - 1)
                start_x = (13.333 - total_w) / 2
                for ki, kpi in enumerate(kpis):
                    x = start_x + ki * (card_w + 0.3)
                    _txt(sl.shapes.add_textbox(Inches(x), Inches(5.6), Inches(card_w), Inches(0.3)),
                         (kpi.get("label") or "").upper(), size=8, color="meta", align=PP_ALIGN.CENTER)
                    _txt(sl.shapes.add_textbox(Inches(x), Inches(5.9), Inches(card_w), Inches(0.5)),
                         kpi.get("value") or "—", size=22, bold=True, align=PP_ALIGN.CENTER)

        elif layout == "agenda":
            _txt(sl.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.3), Inches(0.8)),
                 title or "Agenda", size=32, bold=True)
            for ai, item in enumerate(agenda):
                if not item:
                    continue
                # Number
                _txt(sl.shapes.add_textbox(Inches(1), Inches(2.2 + ai * 0.7), Inches(0.8), Inches(0.5)),
                     f"{ai + 1:02d}", size=20, bold=True, color="accent")
                # Text
                _txt(sl.shapes.add_textbox(Inches(1.9), Inches(2.2 + ai * 0.7), Inches(10), Inches(0.5)),
                     item, size=16)

        elif layout == "comparison":
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=24, bold=True)
            if len(columns) >= 1 and columns[0]:
                _multiline(sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.9), Inches(5.5)),
                           columns[0], size=12)
            if len(columns) >= 2 and columns[1]:
                _multiline(sl.shapes.add_textbox(Inches(6.9), Inches(1.2), Inches(5.9), Inches(5.5)),
                           columns[1], size=12)

        elif layout == "closing":
            _txt(sl.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.0)),
                 title or "Thank You", size=36, bold=True, align=PP_ALIGN.CENTER)
            if subtitle_text:
                _txt(sl.shapes.add_textbox(Inches(1), Inches(3.4), Inches(11.3), Inches(0.6)),
                     subtitle_text, size=16, color="meta", align=PP_ALIGN.CENTER)
            if narrative:
                _multiline(sl.shapes.add_textbox(Inches(3), Inches(4.2), Inches(7.3), Inches(2.0)),
                           narrative, size=12, color="meta")

        else:
            # Fallback: same as chart_full
            _txt(sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5)),
                 title, size=22, bold=True)
            _add_img(sl, _img(vi, 0), 0.3, 0.9, 12.7, 6.3)

    prs.save(buffer)
    buffer.seek(0)
    return buffer


def _generate_slides_pdf(
    slides: List[Dict[str, Any]],
    theme: str = "light",
    report_title: str = "Investment-X Report",
    subtitle: str = "",
    author: str = "",
    classification: str = "For Internal Use Only",
    report_date: Optional[str] = None,
    disclaimer: Optional[str] = None,
) -> BytesIO:
    """Generate a WYSIWYG PDF matching the report editor's light theme appearance."""
    from ix.api.routers.charts.code_execution import render_chart_image

    _require_pdf_dep()
    buffer = BytesIO()

    # Always render charts in light theme for the PDF
    chart_theme = "light"

    # Render all chart images in parallel (multi-slot)
    rendered: Dict[Tuple[int, int], Optional[bytes]] = {}
    figures_to_render: List[Tuple[int, int, dict]] = []
    for i, s in enumerate(slides):
        for slot, key in enumerate(["figure", "figure2", "figure3"]):
            fig = s.get(key)
            if fig and isinstance(fig, dict):
                figures_to_render.append((i, slot, fig))

    if figures_to_render:
        with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
            fmap = {
                executor.submit(render_chart_image, fig, chart_theme): (idx, slot)
                for idx, slot, fig in figures_to_render
            }
            try:
                for future in concurrent.futures.as_completed(fmap, timeout=600):
                    key = fmap[future]
                    try:
                        rendered[key] = future.result()
                    except Exception as exc:
                        logger.error("PDF slide render failed %s: %s", key, exc)
                        rendered[key] = None
            except concurrent.futures.TimeoutError:
                logger.error("PDF slide rendering timed out")

    def _img_tag(i, slot=0, width="269mm"):
        img_bytes = rendered.get((i, slot))
        if not img_bytes:
            return ""
        b64 = base64.b64encode(img_bytes).decode("utf-8")
        return f'<img src="data:image/png;base64,{b64}" style="width:{width};" />'

    # ── Light theme colors (matching editor CSS vars) ──
    bg = "#ffffff"           # slide card background
    fg = "#0a0c14"           # --foreground (dark text)
    muted = "#8b8b8b"        # --muted-foreground
    muted_subtle = "#b0b0b0" # lighter muted
    border = "#e5e5e0"       # --border
    border_light = "#f0efec" # --border/30
    accent = "#3250d2"       # --primary light
    surface = "#f8f7f4"      # --background (page bg)
    card = "#fffefd"         # --card
    success = "#22c55e"
    destructive = "#ef4444"

    date_str = report_date or datetime.now().strftime("%Y-%m-%d")
    esc_title = html_mod.escape(report_title)
    esc_subtitle = html_mod.escape(subtitle) if subtitle else ""
    esc_author = html_mod.escape(author) if author else ""
    esc_classification = html_mod.escape(classification)

    def _header():
        return (
            f'<table style="width:100%; margin-bottom:4mm;">'
            f'<tr>'
            f'<td style="font-size:6.5pt; font-weight:bold; color:{fg}; letter-spacing:2px;">INVESTMENT-X</td>'
            f'<td style="font-size:6.5pt; color:{muted_subtle}; text-align:right;">{date_str}</td>'
            f'</tr></table>'
        )

    def _footer(page_num=""):
        return (
            f'<table style="width:100%; border-top:0.3pt solid {border_light}; '
            f'padding-top:1.5mm; margin-top:auto;">'
            f'<tr>'
            f'<td style="font-size:5.5pt; color:{muted_subtle}; letter-spacing:0.5px;">{esc_classification}</td>'
            f'<td style="font-size:5.5pt; color:{muted_subtle}; text-align:right;">{page_num}</td>'
            f'</tr></table>'
        )

    def _slide_title(title_text):
        """Render a slide title with accent underline, matching the editor h2."""
        if not title_text:
            return ""
        return (
            f'<div style="font-size:15pt; font-weight:bold; color:{fg}; '
            f'padding-bottom:1.5mm; margin-bottom:3mm; '
            f'border-bottom:1.2pt solid {accent};">{title_text}</div>'
        )

    def _kpi_cell(kpi, large=True):
        """Render a single KPI as a bordered card cell."""
        label = html_mod.escape(kpi.get("label", ""))
        value = html_mod.escape(kpi.get("value", "\u2014"))
        change = kpi.get("change", "")
        direction = kpi.get("direction", "")
        change_color = (
            success if direction == "up" else
            destructive if direction == "down" else
            muted_subtle
        )
        val_size = "20pt" if large else "15pt"
        parts_kpi = [
            f'<td style="text-align:center; vertical-align:top; padding:3mm 4mm; '
            f'border:0.4pt solid {border};">',
            f'<div style="font-size:6.5pt; font-weight:bold; color:{muted}; '
            f'letter-spacing:1.2px; text-transform:uppercase; margin-bottom:1.5mm;">{label}</div>',
            f'<div style="font-size:{val_size}; font-weight:bold; color:{fg}; line-height:1;">{value}</div>',
        ]
        if change:
            parts_kpi.append(
                f'<div style="font-size:8pt; color:{change_color}; margin-top:1mm;">'
                f'{html_mod.escape(change)}</div>'
            )
        parts_kpi.append('</td>')
        return "".join(parts_kpi)

    disclaimer_text = _DEFAULT_DISCLAIMER if disclaimer is None else disclaimer

    parts = [f"""<!DOCTYPE html>
<html><head><meta charset="utf-8" />
<style>
    @page {{ size: 297mm 167mm; margin: 8mm 12mm; }}
    body {{ font-family: 'Helvetica', 'Arial', sans-serif; color: {fg}; line-height: 1.45; font-size: 10pt; margin: 0; padding: 0; background: {bg}; }}
    .slide {{ page-break-before: always; page-break-inside: avoid; }}
    .narrative {{ color: {fg}; font-size: 10pt; line-height: 1.65; }}
</style></head><body>

<!-- COVER -->
<div style="padding-top:28mm; text-align:center;">
    <div style="font-size:8pt; font-weight:bold; color:{muted}; letter-spacing:3px; margin-bottom:8mm;">INVESTMENT-X</div>
    <div style="font-size:30pt; font-weight:bold; color:{fg}; margin-bottom:4mm;">{esc_title}</div>"""]

    if esc_subtitle:
        parts.append(f'    <div style="font-size:12pt; color:{muted}; margin-bottom:25mm;">{esc_subtitle}</div>')
    else:
        parts.append('    <div style="margin-bottom:25mm;"></div>')

    meta_items = []
    if esc_author:
        meta_items.append(f'<td style="text-align:center; padding:0 6mm;">'
                          f'<div style="font-size:6pt; font-weight:bold; color:{muted}; letter-spacing:1px; text-transform:uppercase; margin-bottom:1mm;">Author</div>'
                          f'<div style="font-size:8.5pt; color:{fg};">{esc_author}</div></td>')
    meta_items.append(f'<td style="text-align:center; padding:0 6mm;">'
                      f'<div style="font-size:6pt; font-weight:bold; color:{muted}; letter-spacing:1px; text-transform:uppercase; margin-bottom:1mm;">Date</div>'
                      f'<div style="font-size:8.5pt; color:{fg};">{date_str}</div></td>')
    meta_items.append(f'<td style="text-align:center; padding:0 6mm;">'
                      f'<div style="font-size:6pt; font-weight:bold; color:{muted}; letter-spacing:1px; text-transform:uppercase; margin-bottom:1mm;">Classification</div>'
                      f'<div style="font-size:8.5pt; color:{fg};">{esc_classification}</div></td>')
    parts.append(f'<table style="width:55%; margin:0 auto; border-top:0.4pt solid {border}; padding-top:4mm;"><tr>{"".join(meta_items)}</tr></table>')
    parts.append('</div>')

    # ── Content slides — WYSIWYG layout matching editor ──
    for i, s in enumerate(slides):
        layout = s.get("layout", "chart_full")
        title = html_mod.escape(s.get("title") or "")
        subtitle_t = html_mod.escape(s.get("subtitle") or "")
        narrative = html_mod.escape(s.get("narrative") or "").replace("\n", "<br/>")
        kpis = s.get("kpis") or []
        agenda = s.get("agendaItems") or s.get("agenda_items") or []
        columns = s.get("columns") or []

        parts.append('<div class="slide">')
        parts.append(_header())

        if layout == "title":
            parts.append(f'<div style="text-align:center; padding-top:32mm;">')
            parts.append(f'<div style="font-size:28pt; font-weight:bold; color:{fg};">{title}</div>')
            if subtitle_t:
                parts.append(f'<div style="font-size:12pt; color:{muted}; margin-top:4mm;">{subtitle_t}</div>')
            parts.append('</div>')

        elif layout == "section":
            parts.append(f'<div style="padding-top:28mm;">')
            parts.append(f'<div style="background:{accent}; height:2mm; width:25mm; margin-bottom:5mm;"></div>')
            parts.append(f'<div style="font-size:28pt; font-weight:bold; color:{fg};">{title}</div>')
            if subtitle_t:
                parts.append(f'<div style="font-size:12pt; color:{muted}; margin-top:3mm;">{subtitle_t}</div>')
            parts.append('</div>')

        elif layout == "chart_full":
            parts.append(_slide_title(title))
            parts.append(_img_tag(i, 0))

        elif layout == "chart_text":
            parts.append(_slide_title(title))
            parts.append('<table style="width:100%;"><tr>')
            parts.append(f'<td style="width:62%; vertical-align:top;">{_img_tag(i, 0, "100%")}</td>')
            parts.append(f'<td style="width:38%; vertical-align:top; padding-left:4mm;"><div class="narrative">{narrative}</div></td>')
            parts.append('</tr></table>')

        elif layout == "text_chart":
            parts.append(_slide_title(title))
            parts.append('<table style="width:100%;"><tr>')
            parts.append(f'<td style="width:38%; vertical-align:top;"><div class="narrative">{narrative}</div></td>')
            parts.append(f'<td style="width:62%; vertical-align:top; padding-left:4mm;">{_img_tag(i, 0, "100%")}</td>')
            parts.append('</tr></table>')

        elif layout == "text_full":
            parts.append(_slide_title(title))
            if narrative:
                parts.append(f'<div class="narrative" style="font-size:11pt;">{narrative}</div>')

        elif layout == "two_charts":
            parts.append(_slide_title(title))
            parts.append('<table style="width:100%;"><tr>')
            parts.append(f'<td style="width:49%; vertical-align:top;">{_img_tag(i, 0, "100%")}</td>')
            parts.append(f'<td style="width:2%;"></td>')
            parts.append(f'<td style="width:49%; vertical-align:top;">{_img_tag(i, 1, "100%")}</td>')
            parts.append('</tr></table>')

        elif layout == "three_charts":
            parts.append(_slide_title(title))
            parts.append('<table style="width:100%;"><tr>')
            parts.append(f'<td style="width:32%; vertical-align:top;">{_img_tag(i, 0, "100%")}</td>')
            parts.append(f'<td style="width:2%;"></td>')
            parts.append(f'<td style="width:32%; vertical-align:top;">{_img_tag(i, 1, "100%")}</td>')
            parts.append(f'<td style="width:2%;"></td>')
            parts.append(f'<td style="width:32%; vertical-align:top;">{_img_tag(i, 2, "100%")}</td>')
            parts.append('</tr></table>')

        elif layout == "kpi_row":
            parts.append(_slide_title(title))
            if kpis:
                parts.append(f'<table style="width:90%; margin:8mm auto;"><tr>')
                for kpi in kpis:
                    parts.append(_kpi_cell(kpi, large=True))
                parts.append('</tr></table>')

        elif layout == "chart_kpi":
            parts.append(_slide_title(title))
            parts.append(_img_tag(i, 0))
            if kpis:
                parts.append(f'<table style="width:90%; margin:3mm auto;"><tr>')
                for kpi in kpis:
                    parts.append(_kpi_cell(kpi, large=False))
                parts.append('</tr></table>')

        elif layout == "agenda":
            parts.append(_slide_title(title))
            parts.append('<table style="margin-top:5mm; width:80%;">')
            for ai, item in enumerate(agenda):
                if item:
                    parts.append(
                        f'<tr>'
                        f'<td style="padding:2.5mm 5mm 2.5mm 0; width:12mm; vertical-align:top;">'
                        f'<span style="font-size:15pt; font-weight:bold; color:{accent};">{ai + 1:02d}</span></td>'
                        f'<td style="padding:2.5mm 0; vertical-align:top; border-bottom:0.3pt solid {border_light};">'
                        f'<span style="font-size:12pt; color:{fg};">{html_mod.escape(item)}</span></td>'
                        f'</tr>'
                    )
            parts.append('</table>')

        elif layout == "comparison":
            parts.append(_slide_title(title))
            parts.append('<table style="width:100%; margin-top:3mm;"><tr>')
            col_a = html_mod.escape(columns[0]).replace("\n", "<br/>") if len(columns) > 0 and columns[0] else ""
            col_b = html_mod.escape(columns[1]).replace("\n", "<br/>") if len(columns) > 1 and columns[1] else ""
            parts.append(
                f'<td style="width:48%; vertical-align:top; border:0.4pt solid {border}; padding:4mm;">'
                f'<div class="narrative">{col_a}</div></td>'
            )
            parts.append(f'<td style="width:4%;"></td>')
            parts.append(
                f'<td style="width:48%; vertical-align:top; border:0.4pt solid {border}; padding:4mm;">'
                f'<div class="narrative">{col_b}</div></td>'
            )
            parts.append('</tr></table>')

        elif layout == "closing":
            parts.append(f'<div style="text-align:center; padding-top:28mm;">')
            parts.append(f'<div style="font-size:24pt; font-weight:bold; color:{fg};">{title or "Thank You"}</div>')
            if subtitle_t:
                parts.append(f'<div style="font-size:12pt; color:{muted}; margin-top:3mm;">{subtitle_t}</div>')
            if narrative:
                parts.append(f'<div style="font-size:10pt; color:{muted}; margin-top:6mm;">{narrative}</div>')
            parts.append('</div>')

        else:
            parts.append(_slide_title(title))
            parts.append(_img_tag(i, 0))

        parts.append(_footer())
        parts.append('</div>')

    # ── Disclaimer ──
    if disclaimer_text:
        esc_disc = html_mod.escape(disclaimer_text).replace("\n\n", "</p><p>").replace("\n", "<br/>")
        parts.append('<div class="slide">')
        parts.append(_header())
        parts.append(_slide_title("Disclaimer &amp; Disclosures"))
        parts.append(f'<div style="font-size:7.5pt; color:{muted}; line-height:1.7;"><p>{esc_disc}</p></div>')
        parts.append(_footer())
        parts.append('</div>')

    parts.append("</body></html>")

    pisa_status = pisa.CreatePDF("".join(parts), dest=buffer)
    if pisa_status.err:
        logger.error("pisa.CreatePDF failed for slides PDF")

    buffer.seek(0)
    return buffer

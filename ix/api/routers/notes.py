from __future__ import annotations

import base64
import uuid
import textwrap
import re
import html
import json
from datetime import datetime
from io import BytesIO
from typing import List, Optional, Any, Dict
from urllib.parse import urljoin, urlparse
from urllib.request import Request, urlopen

from fastapi import APIRouter, Depends, File, HTTPException, Response, UploadFile, Query
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field, field_validator
from sqlalchemy import func, or_, String
from sqlalchemy.orm import Session

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib import utils
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt

from ix.api.dependencies import get_current_user, get_db
from ix.db.models import InvestmentNote, User

router = APIRouter()

MAX_IMAGE_BYTES = 8 * 1024 * 1024  # 8 MB


class NoteBlock(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    type: str  # text, image, chart, link
    value: Optional[str] = None  # For text
    data: Optional[str] = None  # Base64 for images
    filename: Optional[str] = None
    content_type: Optional[str] = None
    url: Optional[str] = None  # For links
    chart_id: Optional[str] = None  # For charts
    metadata: Dict[str, Any] = Field(default_factory=dict)


class InvestmentNoteCreate(BaseModel):
    title: str = "Untitled Note"
    body: List[Dict[str, Any]] = Field(
        default_factory=lambda: [{"type": "text", "value": "", "id": str(uuid.uuid4())}]
    )
    pinned: bool = False


class InvestmentNoteUpdate(BaseModel):
    title: Optional[str] = None
    body: Optional[List[Dict[str, Any]]] = None
    pinned: Optional[bool] = None


class InvestmentNoteSummary(BaseModel):
    id: str
    title: str
    pinned: bool
    image_count: int = 0
    created_at: datetime
    updated_at: datetime

    class Config:
        from_attributes = True

    @field_validator("created_at", "updated_at", mode="after")
    @classmethod
    def force_utc(cls, v):
        from datetime import timezone

        if v and v.tzinfo is None:
            return v.replace(tzinfo=timezone.utc)
        return v


class InvestmentNoteResponse(BaseModel):
    id: str
    user_id: str
    title: str
    body: List[Dict[str, Any]]
    pinned: bool
    version: int
    created_at: datetime
    updated_at: datetime

    class Config:
        from_attributes = True

    @field_validator("created_at", "updated_at", mode="after")
    @classmethod
    def force_utc(cls, v):
        from datetime import timezone

        if v and v.tzinfo is None:
            return v.replace(tzinfo=timezone.utc)
        return v


class LinkPreviewResponse(BaseModel):
    url: str
    kind: str = "link"
    provider: Optional[str] = None
    title: Optional[str] = None
    subtitle: Optional[str] = None
    description: Optional[str] = None
    image_url: Optional[str] = None


_YOUTUBE_ID_RE = re.compile(r"(?:v=|/shorts/|youtu\.be/)([A-Za-z0-9_-]{11})")


def _extract_youtube_video_id(url: str) -> str:
    value = (url or "").strip()
    if not value:
        return ""
    if re.fullmatch(r"[A-Za-z0-9_-]{11}", value):
        return value
    match = _YOUTUBE_ID_RE.search(value)
    return match.group(1) if match else ""


def _extract_meta_content(raw_html: str, meta_names: list[str]) -> Optional[str]:
    for name in meta_names:
        patterns = [
            rf'<meta[^>]+(?:property|name)=["\']{re.escape(name)}["\'][^>]+content=["\']([^"\']+)["\']',
            rf'<meta[^>]+content=["\']([^"\']+)["\'][^>]+(?:property|name)=["\']{re.escape(name)}["\']',
        ]
        for pattern in patterns:
            match = re.search(pattern, raw_html, flags=re.IGNORECASE)
            if match:
                value = html.unescape(match.group(1).strip())
                if value:
                    return value
    return None


def _fetch_generic_link_preview(url: str) -> LinkPreviewResponse:
    request = Request(url, headers={"User-Agent": "investment-x/1.0"})
    with urlopen(request, timeout=10) as response:
        final_url = response.geturl()
        raw_html = response.read(512_000).decode("utf-8", errors="ignore")

    title_match = re.search(
        r"<title[^>]*>(.*?)</title>", raw_html, flags=re.IGNORECASE | re.DOTALL
    )
    title = _extract_meta_content(raw_html, ["og:title", "twitter:title"]) or (
        html.unescape(title_match.group(1).strip()) if title_match else None
    )
    description = _extract_meta_content(
        raw_html, ["og:description", "description", "twitter:description"]
    )
    image_url = _extract_meta_content(raw_html, ["og:image", "twitter:image"])
    provider = _extract_meta_content(raw_html, ["og:site_name", "application-name"])

    if image_url:
        image_url = urljoin(final_url, image_url)

    return LinkPreviewResponse(
        url=final_url,
        kind="link",
        provider=provider or urlparse(final_url).hostname,
        title=title or final_url,
        description=description,
        image_url=image_url,
    )


def _fetch_youtube_link_preview(url: str) -> Optional[LinkPreviewResponse]:
    video_id = _extract_youtube_video_id(url)
    if not video_id:
        return None

    watch_url = f"https://www.youtube.com/watch?v={video_id}"
    oembed_url = f"https://www.youtube.com/oembed?url={watch_url}&format=json"
    request = Request(oembed_url, headers={"User-Agent": "investment-x/1.0"})

    with urlopen(request, timeout=12) as response:
        payload = json.loads(response.read().decode("utf-8"))

    thumbnail_url = (
        str(payload.get("thumbnail_url") or "").strip()
        or f"https://i.ytimg.com/vi/{video_id}/hqdefault.jpg"
    )

    return LinkPreviewResponse(
        url=watch_url,
        kind="youtube",
        provider="YouTube",
        title=str(payload.get("title") or f"Video {video_id}").strip(),
        subtitle=str(payload.get("author_name") or "YouTube").strip(),
        image_url=thumbnail_url,
    )


def _assert_note_owner(note: InvestmentNote, user: User) -> None:
    if str(note.user_id) != str(user.id):
        raise HTTPException(status_code=404, detail="Note not found.")


def _get_owned_note(db: Session, note_id: str, user: User) -> InvestmentNote:
    note = db.query(InvestmentNote).filter(InvestmentNote.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found.")
    _assert_note_owner(note, user)
    return note


@router.get("/notes", response_model=List[InvestmentNoteSummary])
def list_notes(
    q: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    query = db.query(InvestmentNote).filter(
        InvestmentNote.user_id == str(current_user.id)
    )

    if q:
        search_text = q.strip()
        query = query.filter(
            or_(
                InvestmentNote.title.ilike(f"%{search_text}%"),
                InvestmentNote.body.cast(String).ilike(f"%{search_text}%"),
            )
        )

    notes = query.order_by(
        InvestmentNote.pinned.desc(), InvestmentNote.updated_at.desc()
    ).all()

    summaries = []
    for note in notes:
        img_count = sum(
            1 for block in (note.body or []) if block.get("type") == "image"
        )
        summaries.append(
            InvestmentNoteSummary(
                id=str(note.id),
                title=note.title or "Untitled Note",
                pinned=bool(note.pinned),
                image_count=img_count,
                created_at=note.created_at,
                updated_at=note.updated_at,
            )
        )
    return summaries


@router.get("/notes/{note_id}", response_model=InvestmentNoteResponse)
def get_note(
    note_id: str,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = _get_owned_note(db, note_id, current_user)
    return InvestmentNoteResponse.model_validate(note)


@router.get("/notes/link-preview", response_model=LinkPreviewResponse)
def get_link_preview(
    url: str = Query(..., min_length=3),
    current_user: User = Depends(get_current_user),
):
    value = (url or "").strip()
    if not value:
        raise HTTPException(status_code=400, detail="URL is required.")
    if not re.match(r"^https?://", value, flags=re.IGNORECASE):
        value = f"https://{value}"

    try:
        youtube_preview = _fetch_youtube_link_preview(value)
        if youtube_preview:
            return youtube_preview
        return _fetch_generic_link_preview(value)
    except Exception as exc:
        raise HTTPException(
            status_code=400, detail=f"Unable to load link preview: {exc}"
        ) from exc


@router.post("/notes", response_model=InvestmentNoteResponse, status_code=201)
def create_note(
    payload: InvestmentNoteCreate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = InvestmentNote(
        user_id=str(current_user.id),
        title=(payload.title or "Untitled Note").strip()[:255] or "Untitled Note",
        body=payload.body,
        pinned=bool(payload.pinned),
    )
    db.add(note)
    db.commit()
    db.refresh(note)
    return InvestmentNoteResponse.model_validate(note)


@router.put("/notes/{note_id}", response_model=InvestmentNoteResponse)
def update_note(
    note_id: str,
    payload: InvestmentNoteUpdate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = _get_owned_note(db, note_id, current_user)
    if payload.title is not None:
        note.title = payload.title.strip()[:255] or "Untitled Note"
    if payload.body is not None:
        note.body = payload.body
    if payload.pinned is not None:
        note.pinned = bool(payload.pinned)

    note.version += 1
    note.updated_at = datetime.utcnow()

    db.add(note)
    db.commit()
    db.refresh(note)
    return InvestmentNoteResponse.model_validate(note)


@router.delete("/notes/{note_id}", status_code=204)
def delete_note(
    note_id: str,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = _get_owned_note(db, note_id, current_user)
    db.delete(note)
    db.commit()
    return Response(status_code=204)


@router.post("/notes/{note_id}/images", response_model=Dict[str, Any], status_code=201)
async def upload_note_image(
    note_id: str,
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = _get_owned_note(db, note_id, current_user)

    content_type = str(file.content_type or "").lower()
    if not content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="Only image files are allowed.")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")
    if len(data) > MAX_IMAGE_BYTES:
        raise HTTPException(status_code=400, detail="Image exceeds 8MB size limit.")

    base64_data = base64.b64encode(data).decode("utf-8")

    image_block = {
        "id": str(uuid.uuid4()),
        "type": "image",
        "data": base64_data,
        "filename": (file.filename or "note-image")[:255],
        "content_type": content_type[:128],
        "created_at": datetime.utcnow().isoformat(),
    }

    current_body = list(note.body or [])
    current_body.append(image_block)
    note.body = current_body

    note.updated_at = datetime.utcnow()
    note.version += 1
    db.add(note)
    db.commit()

    return {"id": image_block["id"], "url": f"/api/notes/images/{image_block['id']}"}


@router.get("/notes/images/{image_id}")
def get_note_image(
    image_id: str,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = (
        db.query(InvestmentNote)
        .filter(
            InvestmentNote.user_id == str(current_user.id),
            InvestmentNote.body.cast(String).contains(image_id),
        )
        .first()
    )

    if not note:
        raise HTTPException(status_code=404, detail="Image not found.")

    image_block = next((b for b in note.body if b.get("id") == image_id), None)
    if not image_block or not image_block.get("data"):
        raise HTTPException(status_code=404, detail="Image data not found.")

    image_bytes = base64.b64decode(image_block["data"])
    safe_name = (image_block.get("filename") or "note-image").replace('"', "")

    return StreamingResponse(
        BytesIO(image_bytes),
        media_type=image_block.get("content_type") or "image/png",
        headers={
            "Content-Disposition": f'inline; filename="{safe_name}"',
            "Cache-Control": "private, max-age=604800",
        },
    )


@router.delete("/notes/images/{image_id}", status_code=204)
def delete_note_image(
    image_id: str,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = (
        db.query(InvestmentNote)
        .filter(
            InvestmentNote.user_id == str(current_user.id),
            InvestmentNote.body.cast(String).contains(image_id),
        )
        .first()
    )

    if not note:
        raise HTTPException(status_code=404, detail="Image not found.")

    new_body = [b for b in note.body if b.get("id") != image_id]
    note.body = new_body
    note.updated_at = datetime.utcnow()
    note.version += 1

    db.add(note)
    db.commit()
    return Response(status_code=204)


def _clean_html_for_export(html: str) -> str:
    """Very basic HTML tag removal for export rendering."""
    if not html:
        return ""
    text = (
        html.replace("<br>", "\n")
        .replace("<br/>", "\n")
        .replace("</p>", "\n\n")
        .replace("</div>", "\n")
        .replace("<li>", "\n • ")
        .replace("</li>", "")
    )
    text = re.sub(r"<[^>]+>", "", text)
    text = (
        text.replace("&nbsp;", " ")
        .replace("&amp;", "&")
        .replace("&lt;", "<")
        .replace("&gt;", ">")
        .replace("&quot;", '"')
    )
    return text.strip()


from xhtml2pdf import pisa
from ix.api.routers.custom import render_chart_image
from ix.db.models import CustomChart


from bs4 import BeautifulSoup


@router.get("/notes/{note_id}/export")
def export_note(
    note_id: str,
    format: str = Query("pdf"),  # "pdf" or "pptx"
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = _get_owned_note(db, note_id, current_user)
    title = note.title or "Untitled Report"

    filename = f"InvestmentX_Report_{datetime.now().strftime('%Y%m%d')}.{format}"

    if format.lower() == "pptx":
        prs = Presentation()

        # Slide 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = (
            f"Report Generated: {datetime.now().strftime('%Y-%m-%d')}\nInvestment-X Intelligence"
        )

        for block in note.body or []:
            if block.get("type") == "text" and block.get("value"):
                content = _clean_html_for_export(block["value"])
                if not content:
                    continue
                lines = content.split("\n")
                chunks = [lines[i : i + 12] for i in range(0, len(lines), 12)]
                for i, chunk in enumerate(chunks):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = title + (
                        f" ({i+1})" if len(chunks) > 1 else ""
                    )
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.word_wrap = True
                    tf.text = "\n".join(chunk)
                    for p in tf.paragraphs:
                        p.font.size = Pt(16)

            elif block.get("type") == "image" and block.get("data"):
                try:
                    img_data = base64.b64decode(block["data"])
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(
                        BytesIO(img_data), Inches(1), Inches(1), width=Inches(8)
                    )
                except:
                    continue

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )

    else:  # PDF via xhtml2pdf to preserve tables natively
        html_parts = []
        html_parts.append(
            f"""
        <html>
        <head>
            <style>
                @page {{ size: letter; margin: 2cm; }}
                body {{ font-family: Helvetica, sans-serif; font-size: 11pt; line-height: 1.5; color: #1e293b; }}
                h1 {{ font-size: 24pt; color: #0f172a; text-align: center; margin-bottom: 5px; }}
                .meta {{ text-align: center; font-size: 11pt; color: #64748b; margin-bottom: 30px; border-bottom: 1px solid #cbd5e1; padding-bottom: 15px; }}
                p {{ margin-bottom: 15px; }}
                img {{ max-width: 100%; max-height: 700px; display: block; margin: 20px auto; }}
                table {{ width: 100%; border-collapse: collapse; margin-bottom: 20px; }}
                th, tx {{ border: 1px solid #cbd5e1; padding: 10px; text-align: left; }}
                th {{ background-color: #f1f5f9; font-weight: bold; color: #334155; }}
                td {{ border: 1px solid #e2e8f0; padding: 8px; color: #475569; }}
                ul, ol {{ margin-bottom: 15px; margin-left: 20px; }}
                li {{ margin-bottom: 5px; }}
                .chart-title {{ text-align: center; color: #64748b; font-size: 10pt; font-style: italic; margin-top: -15px; margin-bottom: 25px; }}
            </style>
        </head>
        <body>
            <h1>{html.escape(title)}</h1>
            <div class="meta">Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}</div>
        """
        )

        for block in note.body or []:
            b_type = block.get("type")

            if b_type == "text" and block.get("value"):
                html_val = block["value"]

                # Parse embedded chart nodes dynamically
                try:
                    soup = BeautifulSoup(html_val, "html.parser")
                    chart_nodes = soup.find_all(
                        "div", attrs={"data-chart-block": "true"}
                    )

                    for node in chart_nodes:
                        chart_id = node.get("data-chart-id")
                        if chart_id:
                            db_chart = (
                                db.query(CustomChart).filter_by(id=chart_id).first()
                            )
                            if db_chart and db_chart.figure:
                                chart_bytes = render_chart_image(
                                    db_chart.figure, theme="light"
                                )
                                if chart_bytes:
                                    b64_chart = base64.b64encode(chart_bytes).decode(
                                        "utf-8"
                                    )
                                    c_title = html.escape(
                                        db_chart.name or "Embedded Chart"
                                    )

                                    new_wrapper = soup.new_tag("div")
                                    img_tag = soup.new_tag(
                                        "img", src=f"data:image/png;base64,{b64_chart}"
                                    )
                                    title_div = soup.new_tag(
                                        "div", attrs={"class": "chart-title"}
                                    )
                                    title_div.string = c_title

                                    new_wrapper.append(img_tag)
                                    new_wrapper.append(title_div)
                                    node.replace_with(new_wrapper)

                    html_parts.append(str(soup))
                except Exception as e:
                    import logging

                    logging.getLogger(__name__).error(f"Failed parsing PDF HTML: {e}")
                    html_parts.append(html_val)

            elif b_type == "image" and block.get("data"):
                try:
                    mime = block.get("content_type", "image/png")
                    img_data = block["data"]
                    html_parts.append(f'<img src="data:{mime};base64,{img_data}" />')
                except Exception:
                    pass

        html_parts.append("</body></html>")
        final_html = "\n".join(html_parts)

        # Generate PDF using pisa (xhtml2pdf)
        buffer = BytesIO()
        pisa_status = pisa.CreatePDF(final_html, dest=buffer)

        if pisa_status.err:
            raise HTTPException(
                status_code=500, detail="Failed to generate PDF document layout."
            )

        buffer.seek(0)
        return StreamingResponse(
            buffer,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )
